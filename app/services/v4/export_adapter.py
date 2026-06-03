"""
Export adapter for WYSIWYG HTML/PDF/PPTX export.

Guarantees editor matches export by using the same design tokens and layout logic.
"""

import io
from typing import List

from app.models.v4 import (
    CompiledSlide,
    FlowNode,
    MetricBlock,
    ResolvedDesignTokens,
    TextBlock,
)

try:
    from pptx import Presentation as PptxPresentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class ExportAdapter:
    """Exports compiled slides to HTML, PDF, or PPTX."""

    def __init__(
        self,
        slides: List[CompiledSlide],
        tokens: ResolvedDesignTokens,
    ):
        """
        Initialize exporter.

        Args:
            slides: List of CompiledSlide to export
            tokens: ResolvedDesignTokens for styling
        """
        self.slides = slides
        self.tokens = tokens

    def to_html(self) -> str:
        """
        Export to complete HTML document.

        Uses CSS variables for design tokens and flexbox/grid for layouts.

        Returns:
            Complete HTML string
        """
        palette = self.tokens.palette
        fonts = self.tokens.fonts

        # Build CSS variables from palette
        css_vars = f"""
        :root {{
            --primary: {palette.primary};
            --secondary: {palette.secondary};
            --accent: {palette.accent};
            --background: {palette.background};
            --surface: {palette.surface};
            --surface-alt: {palette.surface_alt};
            --text-primary: {palette.text_primary};
            --text-secondary: {palette.text_secondary};
            --text-muted: {palette.text_muted};
            --border: {palette.border};
            --gradient-start: {palette.gradient_start};
            --gradient-end: {palette.gradient_end};
            --success: {palette.success};
            --warning: {palette.warning};
            --danger: {palette.danger};
            --chart: {palette.chart};
            
            --font-heading: {fonts.heading};
            --font-body: {fonts.body};
            --font-display: {fonts.display};
            --font-mono: {fonts.mono};
        }}
        """

        # Base styles
        styles = """
        * { margin: 0; padding: 0; box-sizing: border-box; }
        html, body { 
            width: 100%; 
            height: 100%; 
            background: var(--background); 
            color: var(--text-primary);
            font-family: var(--font-body);
        }
        
        .slide {
            width: 1920px;
            height: 1080px;
            page-break-after: always;
            position: relative;
            overflow: hidden;
            background: var(--surface);
            display: flex;
            flex-direction: column;
        }
        
        .slide-header {
            padding: 60px;
            border-bottom: 1px solid var(--border);
        }
        
        .slide h1 {
            font-family: var(--font-heading);
            font-size: 56px;
            font-weight: bold;
            color: var(--text-primary);
            margin-bottom: 12px;
        }
        
        .slide h2 {
            font-family: var(--font-heading);
            font-size: 36px;
            font-weight: 600;
            color: var(--text-secondary);
        }
        
        .slide p {
            font-size: 18px;
            line-height: 1.6;
            color: var(--text-primary);
        }
        
        /* Layout: Hero */
        .layout-hero {
            justify-content: center;
            align-items: center;
            text-align: center;
            padding: 120px;
        }
        
        .layout-hero h1 {
            font-size: 72px;
            margin-bottom: 24px;
        }
        
        /* Layout: Process Flow */
        .layout-process-flow {
            padding: 60px;
            display: flex;
            flex-direction: column;
        }
        
        .process-nodes {
            flex: 1;
            display: flex;
            align-items: center;
            justify-content: center;
            gap: 32px;
            flex-wrap: wrap;
            padding: 60px 0;
        }
        
        .process-node {
            background: var(--surface-alt);
            border: 2px solid var(--border);
            border-radius: 12px;
            padding: 24px;
            min-width: 200px;
            max-width: 320px;
            text-align: center;
            position: relative;
        }
        
        .process-node-label {
            font-weight: bold;
            font-size: 18px;
            color: var(--text-primary);
            margin-bottom: 8px;
        }
        
        .process-node-description {
            font-size: 14px;
            color: var(--text-secondary);
        }
        
        .process-connector {
            width: 48px;
            height: 2px;
            background: var(--primary);
            opacity: 0.4;
        }
        
        /* Layout: Metrics */
        .layout-metrics {
            padding: 60px;
            display: grid;
            grid-template-columns: repeat(2, 1fr);
            gap: 24px;
            flex: 1;
        }
        
        .metric-item {
            background: var(--surface-alt);
            border: 1px solid var(--border);
            border-radius: 12px;
            padding: 32px;
            text-align: center;
        }
        
        .metric-value {
            font-family: var(--font-display);
            font-size: 48px;
            font-weight: bold;
            color: var(--primary);
            margin-bottom: 12px;
        }
        
        .metric-label {
            font-size: 16px;
            color: var(--text-secondary);
        }
        
        .metric-delta {
            font-size: 14px;
            margin-top: 8px;
            color: var(--success);
        }
        
        /* Layout: Bento */
        .layout-bento {
            padding: 60px;
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 16px;
            flex: 1;
        }
        
        .bento-item {
            background: var(--surface-alt);
            border: 1px solid var(--border);
            border-radius: 12px;
            padding: 24px;
        }
        
        /* Layout: Split */
        .layout-split {
            display: grid;
            grid-template-columns: 1fr 1fr;
            align-items: center;
            height: 100%;
        }
        
        .split-left, .split-right {
            padding: 60px;
        }
        
        .split-right {
            background: var(--surface-alt);
            border-left: 1px solid var(--border);
        }
        
        /* Layout: Stat Hero */
        .layout-stat-hero {
            justify-content: center;
            align-items: center;
            text-align: center;
            padding: 120px;
        }
        
        .stat-value {
            font-family: var(--font-display);
            font-size: 96px;
            font-weight: bold;
            color: var(--primary);
            line-height: 1;
            margin: 32px 0;
        }
        
        .stat-label {
            font-size: 36px;
            color: var(--text-secondary);
        }
        
        @media print {
            body { margin: 0; }
            .slide { page-break-after: always; }
        }
        """

        # Build HTML slides
        slides_html = ""
        for slide in self.slides:
            slides_html += self._render_slide_html(slide)

        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Presentation Export</title>
    <style>
        {css_vars}
        {styles}
    </style>
</head>
<body>
    {slides_html}
</body>
</html>"""

        return html

    def _render_slide_html(self, slide: CompiledSlide) -> str:
        """Render a single slide to HTML."""
        layout = slide.layout_type
        headline = slide.content.headline or ""
        subhead = slide.content.subhead or ""

        header = ""
        if headline:
            header = f'<div class="slide-header"><h1>{headline}</h1>'
            if subhead:
                header += f'<h2>{subhead}</h2>'
            header += "</div>"

        content = ""
        if layout == "hero":
            content = f'<div class="layout-hero">{header}</div>'
        elif layout == "stat_hero":
            metric = next(
                (b for b in (slide.content.body_blocks or [])
                 if isinstance(b, MetricBlock)),
                None,
            )
            content = f'''<div class="layout-stat-hero">
                <div class="stat-label">{subhead or metric.label if metric else ""}</div>
                <div class="stat-value">{metric.value if metric else "TBD"}</div>
            </div>'''
        elif layout == "metrics":
            metrics_html = ""
            for block in slide.content.body_blocks or []:
                if isinstance(block, MetricBlock):
                    delta_html = (
                        f'<div class="metric-delta">{block.delta_direction}: {block.delta}</div>'
                        if block.delta
                        else ""
                    )
                    metrics_html += f'''<div class="metric-item">
                        <div class="metric-value">{block.value}</div>
                        <div class="metric-label">{block.label}</div>
                        {delta_html}
                    </div>'''
            content = f'<div class="layout-metrics">{header}{metrics_html}</div>'
        elif layout == "process_flow":
            nodes_html = ""
            for i, node in enumerate(slide.content.nodes or []):
                nodes_html += f'''<div class="process-node">
                    <div class="process-node-label">{node.label}</div>
                    <div class="process-node-description">{node.description or ""}</div>
                </div>'''
                if i < len(slide.content.nodes or []) - 1:
                    nodes_html += '<div class="process-connector"></div>'

            content = f'''<div class="layout-process-flow">
                {header}
                <div class="process-nodes">{nodes_html}</div>
            </div>'''
        else:  # default/split
            content = f'<div class="layout-split">{header}</div>'

        return f'<section class="slide">{content}</section>'

    def to_pptx(self) -> bytes:
        """
        Export to PowerPoint format.

        Returns:
            Bytes of PPTX file
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")

        prs = PptxPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Convert palette colors
        bg_color = self._hex_to_rgb(self.tokens.palette.background)
        surface_color = self._hex_to_rgb(self.tokens.palette.surface)
        surface_alt_color = self._hex_to_rgb(self.tokens.palette.surface_alt)
        text_color = self._hex_to_rgb(self.tokens.palette.text_primary)
        primary_color = self._hex_to_rgb(self.tokens.palette.primary)
        border_color = self._hex_to_rgb(self.tokens.palette.border)

        for slide_data in self.slides:
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_color)

            # Add title
            if slide_data.content.headline:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(12), Inches(1.5)
                )
                title_frame = title_box.text_frame
                title_frame.text = slide_data.content.headline
                title_frame.paragraphs[0].font.size = Pt(44)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_color)

            # Render layout-specific content
            if slide_data.layout_type == "process_flow":
                self._render_process_flow_pptx(
                    slide, slide_data, Pt, Inches, RGBColor,
                    surface_alt_color, border_color, text_color
                )
            elif slide_data.layout_type == "metrics":
                self._render_metrics_pptx(
                    slide, slide_data, Pt, Inches, RGBColor,
                    surface_alt_color, text_color, primary_color
                )

        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()

    def _render_process_flow_pptx(
        self, slide: Any, slide_data: CompiledSlide,
        Pt: Any, Inches: Any, RGBColor: Any,
        surface_alt_color, border_color, text_color
    ) -> None:
        """Render process_flow layout in PPTX."""
        nodes = slide_data.content.nodes or []
        if not nodes:
            return

        node_width = 2.5
        gap = 0.5
        total_width = len(nodes) * node_width + (len(nodes) - 1) * gap
        start_x = (13.333 - total_width) / 2

        for i, node in enumerate(nodes):
            x = start_x + i * (node_width + gap)
            y = 2.5

            # Add node rectangle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(node_width),
                Inches(1.0),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*surface_alt_color)
            shape.line.color.rgb = RGBColor(*border_color)

            # Add text
            text_frame = shape.text_frame
            text_frame.text = node.label
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(*text_color)
            text_frame.word_wrap = True

            # Add connector line
            if i < len(nodes) - 1:
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    int(Inches(x + node_width)),
                    int(Inches(y + 0.5)),
                    int(Inches(x + node_width + gap)),
                    int(Inches(y + 0.5)),
                )
                connector.line.color.rgb = RGBColor(*border_color)

    def _render_metrics_pptx(
        self, slide: Any, slide_data: CompiledSlide,
        Pt: Any, Inches: Any, RGBColor: Any,
        surface_alt_color, text_color, primary_color
    ) -> None:
        """Render metrics layout in PPTX."""
        metrics = [
            b for b in (slide_data.content.body_blocks or [])
            if isinstance(b, MetricBlock)
        ]

        cols = 2
        for idx, metric in enumerate(metrics[:4]):  # Max 4 metrics
            col = idx % cols
            row = idx // cols

            x = 0.5 + col * 6.2
            y = 2.5 + row * 2.0

            # Metric box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x),
                Inches(y),
                Inches(6),
                Inches(1.8),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*surface_alt_color)
            box.line.color.rgb = RGBColor(*text_color)

            # Value
            value_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(5.6), Inches(0.8))
            value_frame = value_box.text_frame
            value_frame.text = metric.value
            value_frame.paragraphs[0].font.size = Pt(32)
            value_frame.paragraphs[0].font.bold = True
            value_frame.paragraphs[0].font.color.rgb = RGBColor(*primary_color)

            # Label
            label_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.0), Inches(5.6), Inches(0.6))
            label_frame = label_box.text_frame
            label_frame.text = metric.label
            label_frame.paragraphs[0].font.size = Pt(14)
            label_frame.paragraphs[0].font.color.rgb = RGBColor(*text_color)

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> tuple:
        """Convert hex color to RGB tuple (0-255)."""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
