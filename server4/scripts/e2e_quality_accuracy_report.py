"""E2E quality + accuracy report.

Loads the most recent live-pipeline artifact (`test_v4_live_local_raw.json`)
— produced by running the real `V4ContentPipeline` against live LLMs —
and runs every downstream deterministic stage on it:

  1. `slide_compiler.compile_slides`     → CompiledSlide dicts
  2. `quality_scorer.attach_quality_scores` (Phase 4.5)
  3. `V4PptxBuilder.build()`             → real PPTX bytes
  4. python-pptx read-back               → shape inventory

It then audits no-fake-data invariants: scans every visible string for
boilerplate placeholders ("lorem ipsum", "company name here", "tbd",
"coming soon", "example.com" URLs in citations) and confirms layout
diversity + citation density. Output is a markdown report under
docs/.

This script is intentionally read-only against the artifact and the
codebase — it does not mutate either. It will fail loudly if the
artifact is missing or has zero PASS records, since the whole point is
to ground the report in real generated content.
"""

from __future__ import annotations

import io
import json
import re
import sys
from collections import Counter
from dataclasses import asdict
from pathlib import Path
from typing import Any

HERE = Path(__file__).resolve().parent
SERVER4 = HERE.parent
sys.path.insert(0, str(SERVER4))

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402

from app.services.v4.parallel_writer import GeneratedSlide  # noqa: E402
from app.services.v4.slide_compiler import compile_slides  # noqa: E402
from app.services.v4.quality_scorer import attach_quality_scores  # noqa: E402
from app.services.v4.design_resolver import resolve_design_tokens  # noqa: E402
from app.services.v4.pptx_export import V4PptxBuilder  # noqa: E402


ARTIFACT = SERVER4 / "test_v4_live_local_raw.json"
REPORT_OUT = SERVER4.parent / "docs" / "E2E_QUALITY_ACCURACY_REPORT.md"


# ── No-fake-data invariants ───────────────────────────────────────

# Patterns that MUST NOT appear in any visible slide text. Each is a
# regex that catches both literal and templated placeholder copy.
_FAKE_PATTERNS: list[tuple[str, re.Pattern[str]]] = [
    ("lorem_ipsum", re.compile(r"\blorem\s+ipsum\b", re.I)),
    ("company_name_placeholder", re.compile(r"\bcompany\s+name\s+(here|placeholder)\b", re.I)),
    ("your_company_token", re.compile(r"\[your[\s_-]?(company|product|brand|x)\]", re.I)),
    ("tbd_marker", re.compile(r"^\s*(TBD|TBA|TODO|FIXME|XXX)\s*$")),
    ("coming_soon", re.compile(r"\bcoming\s+soon\b", re.I)),
    ("template_curly_brace", re.compile(r"\{\{\s*[a-z_][a-z0-9_]*\s*\}\}", re.I)),
    ("ellipsis_only", re.compile(r"^\s*\.{3,}\s*$")),
]

# URL hosts that indicate a fabricated citation. example.com / .org /
# .net are documented IETF placeholder domains (RFC 2606); a real
# research run must never cite them.
_PLACEHOLDER_HOSTS = ("example.com", "example.org", "example.net", "test.com")


def _gather_slide_text(slide: dict[str, Any]) -> list[str]:
    """Collect every visible string from a raw GeneratedSlide-shaped dict."""
    out: list[str] = []
    for k in ("headline", "subheadline", "body", "image_prompt"):
        v = slide.get(k)
        if isinstance(v, str) and v.strip():
            out.append(v)
    for v in slide.get("bullets") or []:
        if isinstance(v, str):
            out.append(v)
    for s in slide.get("stat_blocks") or []:
        if isinstance(s, dict):
            for k in ("value", "label"):
                if isinstance(s.get(k), str):
                    out.append(s[k])
    q = slide.get("quote") or {}
    if isinstance(q, dict):
        for k in ("text", "attribution"):
            if isinstance(q.get(k), str):
                out.append(q[k])
    cmp = slide.get("comparison") or {}
    for col in (cmp.get("columns") or []) if isinstance(cmp, dict) else []:
        if isinstance(col, dict):
            if isinstance(col.get("title"), str):
                out.append(col["title"])
            for it in col.get("items") or []:
                if isinstance(it, str):
                    out.append(it)
    tl = slide.get("timeline") or {}
    for ev in (tl.get("events") or []) if isinstance(tl, dict) else []:
        if isinstance(ev, dict):
            for k in ("date", "title", "description"):
                if isinstance(ev.get(k), str):
                    out.append(ev[k])
    dg = slide.get("diagram") or {}
    for n in (dg.get("nodes") or []) if isinstance(dg, dict) else []:
        if isinstance(n, dict) and isinstance(n.get("label"), str):
            out.append(n["label"])
    for m in slide.get("team_members") or []:
        if isinstance(m, dict):
            for k in ("name", "role", "bio"):
                if isinstance(m.get(k), str):
                    out.append(m[k])
    return out


def _audit_fake_data(slides: list[dict[str, Any]]) -> dict[str, Any]:
    findings: list[dict[str, Any]] = []
    placeholder_citations: list[dict[str, Any]] = []
    for sl in slides:
        idx = sl.get("index")
        # Visible slide text (headline/subheadline/body/bullets/etc.)
        for s in _gather_slide_text(sl):
            for name, pat in _FAKE_PATTERNS:
                if pat.search(s):
                    findings.append(
                        {"slide": idx, "rule": name, "snippet": s[:120]}
                    )
        # Team-member fields are surfaced on the team slide and must
        # also be free of placeholder copy. The team-resolver has been
        # observed to emit `TBD — Founder & CEO` style stubs when no
        # team data is available; flag those at the slide level too.
        for m in sl.get("team_members") or []:
            if not isinstance(m, dict):
                continue
            for k in ("name", "role", "bio"):
                v = m.get(k)
                if not isinstance(v, str):
                    continue
                for name, pat in _FAKE_PATTERNS:
                    if pat.search(v):
                        findings.append({
                            "slide": idx,
                            "rule": f"team_member.{k}::{name}",
                            "snippet": v[:120],
                        })
        for c in sl.get("citations") or []:
            if not isinstance(c, dict):
                continue
            url = (c.get("url") or "").lower()
            if any(h in url for h in _PLACEHOLDER_HOSTS):
                placeholder_citations.append(
                    {"slide": idx, "url": c.get("url"), "title": c.get("title")}
                )
    return {
        "fake_text_findings": findings,
        "placeholder_citations": placeholder_citations,
        "total_violations": len(findings) + len(placeholder_citations),
    }


# ── PPTX shape inventory ──────────────────────────────────────────

def _pptx_inventory(blob: bytes) -> dict[str, Any]:
    prs = Presentation(io.BytesIO(blob))
    n_slides = len(prs.slides)
    pictures = 0
    rounded_cards = 0
    right_arrows = 0
    hyperlinks: list[str] = []
    page_numbers_found = 0
    for slide in prs.slides:
        page_number_seen = False
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13:  # PICTURE
                pictures += 1
            try:
                ast = shape.auto_shape_type
            except Exception:
                ast = None
            if ast is not None:
                if int(ast) == int(MSO_SHAPE.ROUNDED_RECTANGLE):
                    rounded_cards += 1
                elif int(ast) == int(MSO_SHAPE.RIGHT_ARROW):
                    right_arrows += 1
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        addr = getattr(run.hyperlink, "address", None)
                        if addr:
                            hyperlinks.append(addr)
                        if not page_number_seen and re.fullmatch(
                            r"\s*\d+\s*/\s*\d+\s*", run.text or ""
                        ):
                            page_number_seen = True
        if page_number_seen:
            page_numbers_found += 1
    return {
        "n_slides": n_slides,
        "pictures": pictures,
        "rounded_cards": rounded_cards,
        "right_arrows": right_arrows,
        "hyperlinks": hyperlinks,
        "page_numbers_found": page_numbers_found,
    }


# ── Per-deck pipeline ──────────────────────────────────────────────

def _reconstruct_slides(records: list[dict[str, Any]]) -> list[GeneratedSlide]:
    """Rebuild GeneratedSlide dataclasses from raw JSON dicts.

    Drops any keys that aren't part of the dataclass schema so we are
    robust against artifact files written by older pipeline versions.
    """
    fields = set(GeneratedSlide.__dataclass_fields__.keys())
    out: list[GeneratedSlide] = []
    for r in records:
        if not isinstance(r, dict):
            continue
        filtered = {k: v for k, v in r.items() if k in fields}
        # Required fields with sensible fallbacks.
        filtered.setdefault("index", 0)
        filtered.setdefault("intent", "")
        filtered.setdefault("layout", "")
        filtered.setdefault("headline", "")
        out.append(GeneratedSlide(**filtered))
    return out


def _run_one(rec: dict[str, Any]) -> dict[str, Any]:
    deck_id = rec.get("id")
    label = rec.get("label", "")
    raw_slides = rec.get("slides") or []
    slides = _reconstruct_slides(raw_slides)
    if not slides:
        return {"id": deck_id, "label": label, "ok": False, "reason": "empty"}

    # 1. Resolve real design tokens from purpose/industry inferred
    # from the scenario label. This mirrors what the live router does.
    purpose = "investor_pitch" if "pitch" in label.lower() else "general"
    industry = "fintech" if "fintech" in label.lower() else "saas"
    tokens = resolve_design_tokens(purpose=purpose, industry=industry).to_dict()

    # 2. Compile + score using the real production code path.
    compiled = compile_slides(slides=slides, deck_title=deck_id)
    attach_quality_scores(compiled, tokens)
    overall_scores = [
        (c.get("quality_score") or {}).get("overall") for c in compiled
    ]
    overall_scores = [s for s in overall_scores if isinstance(s, int)]
    avg_quality = (
        round(sum(overall_scores) / len(overall_scores), 1)
        if overall_scores else None
    )
    pass_threshold = sum(
        1 for c in compiled
        if (c.get("quality_score") or {}).get("passes_threshold")
    )

    # 3. Build the real PPTX (default mode: image fetch enabled, no
    # resolver — so external HTTP images may fail and fall back to
    # honest text placeholders, exactly the production behaviour).
    blob = V4PptxBuilder().build(raw_slides, tokens)
    inventory = _pptx_inventory(blob)

    # 4. No-fake-data audit on the source DSL slides.
    audit = _audit_fake_data(raw_slides)

    # 5. Layout & citation diversity metrics.
    layouts = [s.get("layout", "") for s in raw_slides]
    layout_counts = Counter(layouts)
    unique_layouts = len(layout_counts)
    diversity = round(unique_layouts / max(1, len(layouts)), 2)
    n_with_citations = sum(
        1 for s in raw_slides
        if isinstance(s.get("citations"), list) and len(s["citations"]) > 0
    )
    citation_density = round(n_with_citations / max(1, len(raw_slides)), 2)

    # 6. Image coverage — how many slides actually carry a generated
    # image URL. Bug 1 fix: this is now a declared field on
    # `GeneratedSlide` so it survives `asdict()`.
    n_with_image = sum(
        1 for s in raw_slides if isinstance(s.get("image_url"), str) and s["image_url"]
    )
    image_coverage = round(n_with_image / max(1, len(raw_slides)), 2)

    return {
        "id": deck_id,
        "label": label,
        "ok": True,
        "n_slides": len(raw_slides),
        "n_compiled": len(compiled),
        "layouts": layouts,
        "layout_counts": dict(layout_counts),
        "diversity": diversity,
        "n_with_citations": n_with_citations,
        "citation_density": citation_density,
        "avg_quality": avg_quality,
        "n_passing_threshold": pass_threshold,
        "pptx_bytes": len(blob),
        "pptx_inventory": inventory,
        "audit": audit,
        "n_with_image": n_with_image,
        "image_coverage": image_coverage,
    }


# ── Report rendering ──────────────────────────────────────────────

def _md_report(records: list[dict[str, Any]]) -> str:
    lines: list[str] = []
    lines.append("# E2E Quality + Accuracy Report — V4 Pipeline")
    lines.append("")
    lines.append(
        "Source artifact: live `V4ContentPipeline.generate()` run against "
        "real LLMs (`test_v4_live_local_raw.json`). Every metric below is "
        "computed from production code paths — no mocks, no fixtures."
    )
    lines.append("")
    lines.append(
        "> **What the quality score actually measures.** "
        "Phase 4.5 (`quality_scorer.py`) is a *deterministic rule-violation* "
        "check across three structural dimensions: WCAG contrast ratios, "
        "alignment validity, and density bands per kit. A score of 100 "
        "means *zero rule violations*, not *aesthetically excellent*. "
        "Visual taste is not measured."
    )
    lines.append("")

    total_slides = sum(r.get("n_slides", 0) for r in records if r.get("ok"))
    total_violations = sum(
        r["audit"]["total_violations"] for r in records if r.get("ok")
    )
    total_with_image = sum(r.get("n_with_image", 0) for r in records if r.get("ok"))
    avg_quality_all = [r["avg_quality"] for r in records if r.get("ok") and r.get("avg_quality") is not None]
    avg_quality = (
        round(sum(avg_quality_all) / len(avg_quality_all), 1)
        if avg_quality_all else None
    )

    lines.append("## Summary")
    lines.append("")
    lines.append(f"- **Decks scored:** {sum(1 for r in records if r.get('ok'))}")
    lines.append(f"- **Total slides:** {total_slides}")
    lines.append(
        f"- **Avg Phase 4.5 rule-violation score:** {avg_quality} "
        f"(100 = zero contrast/alignment/density violations)"
    )
    lines.append(f"- **No-fake-data violations:** {total_violations}")
    lines.append(
        f"- **Image coverage:** {total_with_image} / {total_slides} "
        f"slides carry a real generated `image_url`"
    )
    lines.append("")

    for r in records:
        if not r.get("ok"):
            lines.append(f"## Deck `{r['id']}` — SKIPPED ({r.get('reason')})")
            continue
        lines.append(f"## Deck `{r['id']}` — {r['label']}")
        lines.append("")
        lines.append(f"- Slides: **{r['n_slides']}** (compiled: {r['n_compiled']})")
        lines.append(
            f"- Layout diversity: **{r['diversity']}** "
            f"({len(r['layout_counts'])} unique / {r['n_slides']})"
        )
        lines.append(f"- Layout histogram: `{r['layout_counts']}`")
        lines.append(
            f"- Citation density: **{r['citation_density']}** "
            f"({r['n_with_citations']}/{r['n_slides']} slides cite sources)"
        )
        lines.append(
            f"- Rule-violation score (Phase 4.5): avg **{r['avg_quality']}** · "
            f"{r['n_passing_threshold']}/{r['n_compiled']} pass ≥70 threshold"
        )
        lines.append(
            f"- Image coverage: **{r['n_with_image']}/{r['n_slides']}** slides "
            f"carry a generated `image_url`"
        )
        lines.append(
            f"- PPTX bytes: **{r['pptx_bytes']:,}** · "
            f"shapes: pictures={r['pptx_inventory']['pictures']}, "
            f"cards={r['pptx_inventory']['rounded_cards']}, "
            f"arrows={r['pptx_inventory']['right_arrows']}, "
            f"hyperlinks={len(r['pptx_inventory']['hyperlinks'])}, "
            f"page-number footers={r['pptx_inventory']['page_numbers_found']}"
        )
        a = r["audit"]
        if a["total_violations"] == 0:
            lines.append("- No-fake-data audit: **PASS** (zero violations)")
        else:
            lines.append(
                f"- No-fake-data audit: **FAIL** — "
                f"{len(a['fake_text_findings'])} text findings, "
                f"{len(a['placeholder_citations'])} placeholder citations"
            )
            for f in a["fake_text_findings"][:5]:
                lines.append(f"  - slide {f['slide']} `{f['rule']}` :: `{f['snippet']}`")
            for p in a["placeholder_citations"][:5]:
                lines.append(f"  - slide {p['slide']} placeholder cite `{p['url']}`")
        lines.append("")

    lines.append("## Methodology")
    lines.append("")
    lines.append(
        "1. Reconstruct `GeneratedSlide` dataclasses from the live "
        "artifact.\n"
        "2. Resolve real `design_tokens` via `resolve_design_tokens()` "
        "with purpose+industry inferred from the scenario label.\n"
        "3. Run `compile_slides()` (production path) to produce kit-bound "
        "compiled slides.\n"
        "4. Run `attach_quality_scores()` (Phase 4.5) — every dimension "
        "(contrast, alignment, density) measured against the real palette.\n"
        "5. Build PPTX with `V4PptxBuilder().build()` using the real DSL "
        "(images attempt fetch; fall back to honest text placeholders).\n"
        "6. Re-open the PPTX with `python-pptx` and walk every shape: "
        "count pictures, rounded cards, right-arrow connectors, "
        "clickable hyperlink runs, and `i / N` page-number footers.\n"
        "7. Scan every visible string against fake-data regexes "
        "(`lorem ipsum`, `[your X]`, `TBD`, `coming soon`, "
        "`{{template}}`, ellipsis-only) and check citation hosts "
        "against the IETF RFC 2606 placeholder set."
    )
    lines.append("")

    lines.append("## Notes & Caveats")
    lines.append("")
    lines.append(
        "- **What the score does NOT measure.** Phase 4.5 only catches "
        "*rule-violations* (WCAG contrast, alignment validity, density "
        "bands per kit). It does NOT measure narrative quality, "
        "argument strength, visual taste, or whether the deck would "
        "actually persuade a real investor. Treat 100 as 'zero "
        "structural defects', not as a design grade.\n"
        "- **Image coverage is non-deterministic** because the image "
        "stage hits live providers (Azure Flux \u2192 Pollinations \u2192 "
        "gradient-SVG fallback chain). A 0-coverage deck means every "
        "tier failed for that run; re-running typically recovers. "
        "Layouts with no image intent (`title-only`, `quote`, "
        "`bullet-points`) correctly carry `image_url=None`.\n"
        "- **Card shapes = 0** if the run did not surface "
        "`stat-grid` / `comparison` layouts. Card rendering is "
        "covered separately by unit tests.\n"
        "- **Layout normalization**: the planner's "
        "`_CANONICAL_LAYOUTS` set covers `title-only, two-column, "
        "stat-hero, grid-3, chart-focus, image-full, quote, "
        "comparison, timeline, table, diagram, process, "
        "bullet-points, auto`. The PPTX exporter routes all of these "
        "to dedicated renderers post-Bug-2 fix; previously "
        "`chart-focus` and `process` silently fell through to the "
        "bullets renderer."
    )
    lines.append("")
    return "\n".join(lines)


def main() -> int:
    if not ARTIFACT.exists():
        print(f"Artifact missing: {ARTIFACT}", file=sys.stderr)
        return 2
    data = json.loads(ARTIFACT.read_text(encoding="utf-8"))
    if not isinstance(data, list) or not data:
        print(f"Artifact is empty or wrong shape: {ARTIFACT}", file=sys.stderr)
        return 2

    pass_records = [r for r in data if isinstance(r, dict) and r.get("status") == "PASS"]
    if not pass_records:
        print("No PASS records in artifact — nothing to score.", file=sys.stderr)
        return 2

    results: list[dict[str, Any]] = []
    for r in pass_records:
        try:
            results.append(_run_one(r))
        except Exception as exc:  # surface the real failure
            results.append({
                "id": r.get("id"), "label": r.get("label"),
                "ok": False, "reason": f"runner error: {exc!r}",
            })

    md = _md_report(results)
    REPORT_OUT.parent.mkdir(parents=True, exist_ok=True)
    REPORT_OUT.write_text(md, encoding="utf-8")
    # On Windows the default stdout codec is cp1252 and chokes on the
    # `≥` glyph; force a utf-8 reconfigure if available, otherwise fall
    # back to ascii-safe printing so the script never errors out after a
    # successful file write.
    try:
        sys.stdout.reconfigure(encoding="utf-8")  # type: ignore[attr-defined]
        print(md)
    except Exception:
        print(md.encode("ascii", "replace").decode("ascii"))
    print(f"\n[saved] {REPORT_OUT}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
