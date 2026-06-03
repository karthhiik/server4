"""
MCP Tool Registry — Declares all 40+ V7 tools organised by agent domain.

Each tool is an async handler with full type hints and Pydantic validation.
Tools are registered with the FastMCP server in ``mcp_server.py``.

Naming convention: ``<domain>_<verb>_<noun>``
"""

from datetime import datetime
from typing import Any, Optional

from bson import ObjectId
from pydantic import BaseModel, Field

import structlog

from app.database import get_db
from app.models.dsl_v2 import (
    LayoutType,
    PresentationDSL,
    SlideDSL,
    SlideType,
    ThemeDSL,
)

logger = structlog.get_logger(__name__)


# ═══════════════════════════════════════════════════════════════════
# REQUEST / RESPONSE SCHEMAS  (used by tool handlers)
# ═══════════════════════════════════════════════════════════════════

class CreatePresentationRequest(BaseModel):
    title: str = Field(..., min_length=1, max_length=200)
    description: Optional[str] = Field(default=None, max_length=2000)
    archetype: Optional[str] = Field(default=None, max_length=50)
    user_id: str = Field(default="system")


class UpdatePresentationRequest(BaseModel):
    presentation_id: str
    title: Optional[str] = Field(default=None, max_length=200)
    description: Optional[str] = Field(default=None, max_length=2000)


class SlideRequest(BaseModel):
    presentation_id: str
    slide_type: SlideType = SlideType.CUSTOM
    layout: LayoutType = LayoutType.CENTER_FOCUS
    index: Optional[int] = None
    content: dict[str, Any] = Field(default_factory=dict)


class ReorderRequest(BaseModel):
    presentation_id: str
    slide_ids: list[str]


class GenerateDSLRequest(BaseModel):
    topic: str = Field(..., min_length=1, max_length=500)
    description: str = Field(default="", max_length=5000)
    audience: str = Field(default="investors", max_length=200)
    slide_count: int = Field(default=10, ge=3, le=50)
    archetype: Optional[str] = None
    language: str = Field(default="en")


class OutlineRequest(BaseModel):
    topic: str = Field(..., min_length=1, max_length=500)
    description: str = Field(default="", max_length=5000)
    slide_count: int = Field(default=10, ge=3, le=50)


class ThemeRequest(BaseModel):
    presentation_id: str
    theme: Optional[dict[str, Any]] = None
    preset: Optional[str] = None


class ResearchRequest(BaseModel):
    topic: str = Field(..., min_length=1, max_length=500)
    depth: str = Field(default="standard", pattern=r"^(quick|standard|deep)$")


class QARequest(BaseModel):
    presentation_id: str
    checks: list[str] = Field(default_factory=lambda: ["layout", "contrast", "content"])


class VFXRequest(BaseModel):
    slide_id: str
    scene_type: str = Field(default="particles", max_length=50)
    config: dict[str, Any] = Field(default_factory=dict)


class PPTXRequest(BaseModel):
    presentation_id: str
    template: Optional[str] = None


class ToolResult(BaseModel):
    """Standardised tool response envelope."""
    success: bool = True
    data: Any = None
    error: Optional[str] = None
    meta: dict[str, Any] = Field(default_factory=dict)


# ═══════════════════════════════════════════════════════════════════
# TOOL HANDLER IMPLEMENTATIONS
# ═══════════════════════════════════════════════════════════════════

# ── Presentation Management ───────────────────────────────────

async def create_presentation(
    title: str,
    description: Optional[str] = None,
    archetype: Optional[str] = None,
    user_id: str = "system",
) -> dict[str, Any]:
    """Create a new presentation document in MongoDB."""
    db = get_db()
    doc_id = str(ObjectId())
    doc = {
        "_id": doc_id,
        "user_id": user_id,
        "title": title,
        "description": description,
        "archetype": archetype,
        "mode": "premium",
        "created_from": "ai",
        "slide_count": 0,
        "generation_state": "idle",
        "generation_progress": 0,
        "generation_message": "",
        "dsl_version": "2.0",
        "created_at": datetime.utcnow(),
        "updated_at": datetime.utcnow(),
    }
    await db.presentations.insert_one(doc)
    logger.info("tool_create_presentation", id=doc_id, title=title)
    return {"id": doc_id, "title": title}


async def get_presentation(presentation_id: str) -> dict[str, Any]:
    """Retrieve a presentation by ID."""
    db = get_db()
    doc = await db.presentations.find_one({"_id": presentation_id})
    if not doc:
        raise ValueError(f"Presentation {presentation_id} not found")
    doc["id"] = str(doc.pop("_id"))
    return doc


async def update_presentation(
    presentation_id: str,
    title: Optional[str] = None,
    description: Optional[str] = None,
) -> dict[str, Any]:
    """Update presentation metadata."""
    db = get_db()
    updates: dict[str, Any] = {"updated_at": datetime.utcnow()}
    if title is not None:
        updates["title"] = title
    if description is not None:
        updates["description"] = description
    result = await db.presentations.find_one_and_update(
        {"_id": presentation_id},
        {"$set": updates},
        return_document=True,
    )
    if not result:
        raise ValueError(f"Presentation {presentation_id} not found")
    result["id"] = str(result.pop("_id"))
    return result


async def delete_presentation(presentation_id: str) -> dict[str, Any]:
    """Delete a presentation and cascade to slides."""
    db = get_db()
    result = await db.presentations.delete_one({"_id": presentation_id})
    if result.deleted_count == 0:
        raise ValueError(f"Presentation {presentation_id} not found")
    await db.slides.delete_many({"presentation_id": presentation_id})
    await db.slide_versions.delete_many({"presentation_id": presentation_id})
    logger.info("tool_delete_presentation", id=presentation_id)
    return {"deleted": presentation_id}


async def list_presentations(
    user_id: str = "system",
    skip: int = 0,
    limit: int = 20,
) -> list[dict[str, Any]]:
    """List presentations for a user."""
    db = get_db()
    cursor = (
        db.presentations.find({"user_id": user_id})
        .sort("created_at", -1)
        .skip(skip)
        .limit(limit)
    )
    docs = await cursor.to_list(limit)
    for d in docs:
        d["id"] = str(d.pop("_id"))
    return docs


async def duplicate_presentation(presentation_id: str) -> dict[str, Any]:
    """Duplicate a presentation with all its slides."""
    db = get_db()
    original = await db.presentations.find_one({"_id": presentation_id})
    if not original:
        raise ValueError(f"Presentation {presentation_id} not found")
    new_id = str(ObjectId())
    new_doc = {
        **original,
        "_id": new_id,
        "title": f"{original['title']} (Copy)",
        "created_at": datetime.utcnow(),
        "updated_at": datetime.utcnow(),
    }
    await db.presentations.insert_one(new_doc)
    slides = await db.slides.find({"presentation_id": presentation_id}).to_list(200)
    for slide in slides:
        await db.slides.insert_one({
            **slide,
            "_id": str(ObjectId()),
            "presentation_id": new_id,
            "created_at": datetime.utcnow(),
            "updated_at": datetime.utcnow(),
        })
    logger.info("tool_duplicate_presentation", original=presentation_id, new=new_id)
    return {"id": new_id, "title": new_doc["title"]}


# ── Slide Operations ──────────────────────────────────────────

async def add_slide(
    presentation_id: str,
    slide_type: str = "custom",
    layout: str = "center-focus",
    index: Optional[int] = None,
    content: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    """Add a slide to a presentation."""
    db = get_db()
    pres = await db.presentations.find_one({"_id": presentation_id})
    if not pres:
        raise ValueError(f"Presentation {presentation_id} not found")
    if index is None:
        index = pres.get("slide_count", 0)
    slide_id = str(ObjectId())
    doc = {
        "_id": slide_id,
        "presentation_id": presentation_id,
        "index": index,
        "type": slide_type,
        "layout": layout,
        "content": content or {},
        "style": {},
        "elements": [],
        "speakerNotes": None,
        "version": 1,
        "created_at": datetime.utcnow(),
        "updated_at": datetime.utcnow(),
    }
    await db.slides.insert_one(doc)
    await db.presentations.update_one(
        {"_id": presentation_id},
        {"$inc": {"slide_count": 1}, "$set": {"updated_at": datetime.utcnow()}},
    )
    logger.info("tool_add_slide", presentation=presentation_id, slide=slide_id)
    return {"id": slide_id, "index": index}


async def update_slide(
    slide_id: str,
    content: Optional[dict[str, Any]] = None,
    layout: Optional[str] = None,
    style: Optional[dict[str, Any]] = None,
    speaker_notes: Optional[str] = None,
) -> dict[str, Any]:
    """Update slide content, layout, or style."""
    db = get_db()
    updates: dict[str, Any] = {"updated_at": datetime.utcnow()}
    if content is not None:
        updates["content"] = content
    if layout is not None:
        updates["layout"] = layout
    if style is not None:
        updates["style"] = style
    if speaker_notes is not None:
        updates["speakerNotes"] = speaker_notes
    result = await db.slides.find_one_and_update(
        {"_id": slide_id}, {"$set": updates}, return_document=True,
    )
    if not result:
        raise ValueError(f"Slide {slide_id} not found")
    result["id"] = str(result.pop("_id"))
    return result


async def delete_slide(slide_id: str) -> dict[str, Any]:
    """Delete a slide and update presentation count."""
    db = get_db()
    slide = await db.slides.find_one({"_id": slide_id})
    if not slide:
        raise ValueError(f"Slide {slide_id} not found")
    await db.slides.delete_one({"_id": slide_id})
    await db.presentations.update_one(
        {"_id": slide["presentation_id"]},
        {"$inc": {"slide_count": -1}, "$set": {"updated_at": datetime.utcnow()}},
    )
    return {"deleted": slide_id}


async def reorder_slides(
    presentation_id: str,
    slide_ids: list[str],
) -> dict[str, Any]:
    """Reorder slides by setting new index values."""
    db = get_db()
    for idx, sid in enumerate(slide_ids):
        await db.slides.update_one(
            {"_id": sid, "presentation_id": presentation_id},
            {"$set": {"index": idx, "updated_at": datetime.utcnow()}},
        )
    return {"reordered": len(slide_ids)}


async def get_slide(slide_id: str) -> dict[str, Any]:
    """Retrieve a single slide."""
    db = get_db()
    doc = await db.slides.find_one({"_id": slide_id})
    if not doc:
        raise ValueError(f"Slide {slide_id} not found")
    doc["id"] = str(doc.pop("_id"))
    return doc


async def get_all_slides(presentation_id: str) -> list[dict[str, Any]]:
    """Get all slides for a presentation, ordered by index."""
    db = get_db()
    cursor = db.slides.find({"presentation_id": presentation_id}).sort("index", 1)
    docs = await cursor.to_list(200)
    for d in docs:
        d["id"] = str(d.pop("_id"))
    return docs


async def duplicate_slide(slide_id: str) -> dict[str, Any]:
    """Duplicate a slide within its presentation."""
    db = get_db()
    original = await db.slides.find_one({"_id": slide_id})
    if not original:
        raise ValueError(f"Slide {slide_id} not found")
    new_id = str(ObjectId())
    new_doc = {
        **original,
        "_id": new_id,
        "index": original["index"] + 1,
        "version": 1,
        "created_at": datetime.utcnow(),
        "updated_at": datetime.utcnow(),
    }
    await db.slides.insert_one(new_doc)
    await db.presentations.update_one(
        {"_id": original["presentation_id"]},
        {"$inc": {"slide_count": 1}},
    )
    return {"id": new_id, "index": new_doc["index"]}


# ── Content Generation ────────────────────────────────────────

async def generate_dsl(
    topic: str,
    description: str = "",
    audience: str = "investors",
    slide_count: int = 10,
    archetype: Optional[str] = None,
    language: str = "en",
) -> dict[str, Any]:
    """Generate a complete PresentationDSL for the given topic.
    Delegates to the LLM model router for actual generation."""
    from app.services.llm.model_router import get_model_router, TaskType

    router = get_model_router()
    prompt = (
        f"Generate a {slide_count}-slide pitch deck DSL (JSON, version 2.0) "
        f"for topic: '{topic}'. Description: {description}. "
        f"Audience: {audience}. Language: {language}. "
        f"Archetype: {archetype or 'problem-solution'}. "
        f"Return valid JSON matching PresentationDSL v2 schema."
    )
    response = await router.complete(
        task_type=TaskType.STRUCTURED_JSON,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=8000,
    )
    logger.info("tool_generate_dsl", topic=topic, slides=slide_count)
    return {"dsl": response.content, "slide_count": slide_count}


async def compile_react(presentation_dsl: dict[str, Any]) -> dict[str, Any]:
    """Compile a PresentationDSL into React JSX component code."""
    dsl = PresentationDSL.model_validate(presentation_dsl)
    components: list[str] = []
    for slide in dsl.slides:
        components.append(
            f"<Slide id='{slide.id}' type='{slide.type.value}' "
            f"layout='{slide.layout.value}' />"
        )
    jsx = "\n".join(components)
    return {"format": "react", "code": jsx, "slide_count": len(dsl.slides)}


async def compile_revealjs(presentation_dsl: dict[str, Any]) -> dict[str, Any]:
    """Compile a PresentationDSL into Reveal.js HTML sections."""
    dsl = PresentationDSL.model_validate(presentation_dsl)
    sections: list[str] = []
    for slide in dsl.slides:
        bg = ""
        if slide.style.background.colors:
            bg = f' data-background-color="{slide.style.background.colors[0]}"'
        sections.append(
            f'<section id="{slide.id}" data-transition="{slide.revealConfig.transition.value}"{bg}>\n'
            f"  <h2>{slide.content.title}</h2>\n"
            f"  {'<p>' + slide.content.subtitle + '</p>' if slide.content.subtitle else ''}\n"
            f"</section>"
        )
    html = "\n".join(sections)
    return {"format": "revealjs", "html": html, "slide_count": len(dsl.slides)}


async def compile_html(presentation_dsl: dict[str, Any]) -> dict[str, Any]:
    """Compile a PresentationDSL into a standalone HTML page."""
    dsl = PresentationDSL.model_validate(presentation_dsl)
    slides_html: list[str] = []
    for slide in dsl.slides:
        slides_html.append(
            f'<div class="slide" data-type="{slide.type.value}">\n'
            f"  <h2>{slide.content.title}</h2>\n"
            f"  <p>{slide.content.body_text or ''}</p>\n"
            f"</div>"
        )
    page = (
        "<!DOCTYPE html>\n<html>\n<head>\n"
        f"  <title>{dsl.presentation.title}</title>\n"
        "</head>\n<body>\n" + "\n".join(slides_html) + "\n</body>\n</html>"
    )
    return {"format": "html", "html": page, "slide_count": len(dsl.slides)}


async def generate_outline(
    topic: str,
    description: str = "",
    slide_count: int = 10,
) -> dict[str, Any]:
    """Generate a presentation outline (slide titles + types) without full DSL."""
    from app.services.llm.model_router import get_model_router, TaskType

    router = get_model_router()
    prompt = (
        f"Create a {slide_count}-slide outline for: '{topic}'. {description}. "
        f"Return JSON array of {{title, type, key_points}}."
    )
    response = await router.complete(
        task_type=TaskType.OUTLINE_PLANNING,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=2000,
    )
    return {"outline": response.content, "slide_count": slide_count}


async def refine_content(
    slide_id: str,
    instruction: str,
) -> dict[str, Any]:
    """Refine a slide's content using an LLM with the given instruction."""
    db = get_db()
    slide = await db.slides.find_one({"_id": slide_id})
    if not slide:
        raise ValueError(f"Slide {slide_id} not found")

    from app.services.llm.model_router import get_model_router, TaskType

    router = get_model_router()
    prompt = (
        f"Refine this slide content per instruction: '{instruction}'.\n"
        f"Current content: {slide.get('content', {})}\n"
        f"Return the improved content as JSON."
    )
    response = await router.complete(
        task_type=TaskType.REFINEMENT,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=2000,
    )
    return {"refined": response.content, "slide_id": slide_id}


# ── Strategy Tools ────────────────────────────────────────────

async def analyze_presentation(presentation_id: str) -> dict[str, Any]:
    """Analyse a presentation's narrative arc, structure, and coherence."""
    db = get_db()
    pres = await db.presentations.find_one({"_id": presentation_id})
    if not pres:
        raise ValueError(f"Presentation {presentation_id} not found")
    slides = await db.slides.find({"presentation_id": presentation_id}).sort("index", 1).to_list(200)

    analysis = {
        "title": pres.get("title"),
        "slide_count": len(slides),
        "has_title_slide": any(s.get("type") == "title-slide" for s in slides),
        "has_closing_slide": any(s.get("type") == "closing-slide" for s in slides),
        "types_used": list({s.get("type", "custom") for s in slides}),
        "layouts_used": list({s.get("layout", "center-focus") for s in slides}),
        "avg_content_length": (
            sum(len(str(s.get("content", {}))) for s in slides) / max(len(slides), 1)
        ),
    }
    return analysis


async def validate_strategy(presentation_id: str) -> dict[str, Any]:
    """Validate that a presentation follows pitch deck best practices."""
    analysis = await analyze_presentation(presentation_id)
    issues: list[str] = []
    if not analysis["has_title_slide"]:
        issues.append("Missing title slide")
    if not analysis["has_closing_slide"]:
        issues.append("Missing closing/CTA slide")
    if analysis["slide_count"] < 5:
        issues.append("Too few slides (recommend 8-15)")
    if analysis["slide_count"] > 25:
        issues.append("Too many slides (recommend 8-15)")

    score = max(0, 100 - len(issues) * 15)
    return {"score": score, "issues": issues, "analysis": analysis}


async def select_archetype(
    topic: str,
    audience: str = "investors",
) -> dict[str, Any]:
    """Recommend a narrative archetype based on topic and audience."""
    archetypes = {
        "investors": "problem-solution",
        "customers": "before-after",
        "internal": "status-report",
        "educational": "journey",
        "sales": "hero-journey",
    }
    recommended = archetypes.get(audience, "problem-solution")
    return {
        "recommended": recommended,
        "alternatives": list(archetypes.values()),
        "topic": topic,
        "audience": audience,
    }


# ── Research Tools ────────────────────────────────────────────

async def research_topic(
    topic: str,
    depth: str = "standard",
) -> dict[str, Any]:
    """Research a topic using available data sources."""
    from app.services.llm.model_router import get_model_router, TaskType

    router = get_model_router()
    prompt = (
        f"Research the following topic for a pitch deck: '{topic}'. "
        f"Depth: {depth}. Provide key statistics, market data, and citations."
    )
    response = await router.complete(
        task_type=TaskType.GENERAL,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=4000,
    )
    return {"topic": topic, "depth": depth, "data": response.content}


async def extract_document(
    document_text: str,
    extraction_type: str = "key_points",
) -> dict[str, Any]:
    """Extract structured data from a document for slide content."""
    from app.services.llm.model_router import get_model_router, TaskType

    router = get_model_router()
    prompt = (
        f"Extract {extraction_type} from this document:\n{document_text[:5000]}\n"
        f"Return as a structured JSON object."
    )
    response = await router.complete(
        task_type=TaskType.STRUCTURED_JSON,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=3000,
    )
    return {"type": extraction_type, "extracted": response.content}


async def search_web(
    query: str,
    max_results: int = 5,
) -> dict[str, Any]:
    """Search the web for information relevant to a presentation topic."""
    logger.info("tool_search_web", query=query)
    return {
        "query": query,
        "results": [],
        "note": "Web search requires Tavily/Exa API key configuration",
    }


async def analyze_data(
    data: dict[str, Any],
    analysis_type: str = "summary",
) -> dict[str, Any]:
    """Analyse structured data and produce insights for slides."""
    from app.services.llm.model_router import get_model_router, TaskType

    router = get_model_router()
    prompt = (
        f"Analyze this data ({analysis_type}):\n{data}\n"
        f"Provide insights suitable for a pitch deck slide."
    )
    response = await router.complete(
        task_type=TaskType.GENERAL,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=2000,
    )
    return {"analysis_type": analysis_type, "insights": response.content}


# ── Design Tools ──────────────────────────────────────────────

async def apply_theme(
    presentation_id: str,
    theme: Optional[dict[str, Any]] = None,
    preset: Optional[str] = None,
) -> dict[str, Any]:
    """Apply a theme to all slides in a presentation."""
    db = get_db()
    if preset:
        theme_doc = await db.themes.find_one({"name": preset})
        if theme_doc:
            theme = theme_doc.get("config", {})
    if theme is None:
        theme = {}
    await db.presentations.update_one(
        {"_id": presentation_id},
        {"$set": {"theme": theme, "updated_at": datetime.utcnow()}},
    )
    return {"applied": True, "presentation_id": presentation_id}


async def generate_theme(
    industry: str = "technology",
    mood: str = "professional",
    variant: str = "dark",
) -> dict[str, Any]:
    """Generate a theme configuration using AI."""
    from app.services.llm.model_router import get_model_router, TaskType

    router = get_model_router()
    prompt = (
        f"Generate a presentation theme for a {industry} company. "
        f"Mood: {mood}. Variant: {variant}. "
        f"Return JSON with: primaryColor, secondaryColor, accentColor, "
        f"fontHeading, fontBody, backgroundColor."
    )
    response = await router.complete(
        task_type=TaskType.STRUCTURED_JSON,
        messages=[{"role": "user", "content": prompt}],
        max_tokens=1000,
    )
    return {"theme": response.content, "industry": industry, "mood": mood}


async def discover_style(
    reference_url: Optional[str] = None,
    keywords: Optional[list[str]] = None,
) -> dict[str, Any]:
    """Discover visual style inspiration based on references or keywords."""
    return {
        "reference_url": reference_url,
        "keywords": keywords or [],
        "suggestions": [
            {"name": "Minimal Dark", "colors": ["#0a0a0a", "#ffffff", "#3b82f6"]},
            {"name": "Gradient Flow", "colors": ["#1a1a2e", "#16213e", "#0f3460"]},
            {"name": "Clean Light", "colors": ["#ffffff", "#f8fafc", "#1e293b"]},
        ],
    }


async def check_contrast(
    foreground: str,
    background: str,
) -> dict[str, Any]:
    """Check WCAG contrast ratio between two colours."""
    def _hex_to_luminance(hex_color: str) -> float:
        hex_color = hex_color.lstrip("#")
        r, g, b = (int(hex_color[i:i+2], 16) / 255.0 for i in (0, 2, 4))
        components = []
        for c in (r, g, b):
            if c <= 0.03928:
                components.append(c / 12.92)
            else:
                components.append(((c + 0.055) / 1.055) ** 2.4)
        return 0.2126 * components[0] + 0.7152 * components[1] + 0.0722 * components[2]

    l1 = _hex_to_luminance(foreground)
    l2 = _hex_to_luminance(background)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    ratio = (lighter + 0.05) / (darker + 0.05)
    return {
        "foreground": foreground,
        "background": background,
        "ratio": round(ratio, 2),
        "aa_normal": ratio >= 4.5,
        "aa_large": ratio >= 3.0,
        "aaa_normal": ratio >= 7.0,
    }


# ── PPTX Tools ────────────────────────────────────────────────

async def create_pptx(
    presentation_id: str,
    template: Optional[str] = None,
) -> dict[str, Any]:
    """Generate a .pptx file from a presentation's slides."""
    db = get_db()
    slides = await db.slides.find({"presentation_id": presentation_id}).sort("index", 1).to_list(200)
    if not slides:
        raise ValueError(f"No slides found for presentation {presentation_id}")

    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    import io

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        slide_layout = prs.slide_layouts[5]  # blank layout
        pptx_slide = prs.slides.add_slide(slide_layout)
        content = slide_data.get("content", {})
        title_text = content.get("title", "")
        if title_text:
            from pptx.util import Inches, Pt
            txBox = pptx_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
            tf = txBox.text_frame
            tf.text = title_text
            tf.paragraphs[0].font.size = Pt(36)
            tf.paragraphs[0].font.bold = True
        body = content.get("body_text") or content.get("subtitle", "")
        if body:
            txBox2 = pptx_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(4))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            tf2.text = body
            tf2.paragraphs[0].font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    size = buf.tell()
    logger.info("tool_create_pptx", presentation=presentation_id, size_bytes=size)
    return {"presentation_id": presentation_id, "size_bytes": size, "slide_count": len(slides)}


async def add_chart(
    presentation_id: str,
    slide_index: int,
    chart_type: str = "bar",
    chart_data: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    """Add a chart element to a slide."""
    db = get_db()
    result = await db.slides.find_one_and_update(
        {"presentation_id": presentation_id, "index": slide_index},
        {
            "$set": {
                "content.chart_data": chart_data or {},
                "content.chart_type": chart_type,
                "updated_at": datetime.utcnow(),
            }
        },
        return_document=True,
    )
    if not result:
        raise ValueError(f"Slide at index {slide_index} not found")
    return {"slide_index": slide_index, "chart_type": chart_type}


async def add_table(
    presentation_id: str,
    slide_index: int,
    headers: list[str],
    rows: list[list[str]],
) -> dict[str, Any]:
    """Add table data to a slide."""
    db = get_db()
    table_data = {"headers": headers, "rows": rows}
    result = await db.slides.find_one_and_update(
        {"presentation_id": presentation_id, "index": slide_index},
        {
            "$set": {
                "content.table_data": table_data,
                "updated_at": datetime.utcnow(),
            }
        },
        return_document=True,
    )
    if not result:
        raise ValueError(f"Slide at index {slide_index} not found")
    return {"slide_index": slide_index, "rows": len(rows)}


async def add_shape(
    presentation_id: str,
    slide_index: int,
    shape_type: str = "rectangle",
    position: Optional[dict[str, float]] = None,
    size: Optional[dict[str, float]] = None,
    style: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    """Add a shape element to a slide."""
    element = {
        "id": f"shape_{ObjectId()}",
        "type": "shape",
        "content": shape_type,
        "position": position or {"x": 0.1, "y": 0.1},
        "size": size or {"width": 0.3, "height": 0.3},
        "style": style or {},
    }
    db = get_db()
    await db.slides.update_one(
        {"presentation_id": presentation_id, "index": slide_index},
        {"$push": {"elements": element}, "$set": {"updated_at": datetime.utcnow()}},
    )
    return {"element_id": element["id"], "shape_type": shape_type}


async def add_image(
    presentation_id: str,
    slide_index: int,
    image_url: str,
    position: Optional[dict[str, float]] = None,
    size: Optional[dict[str, float]] = None,
    alt_text: Optional[str] = None,
) -> dict[str, Any]:
    """Add an image element to a slide."""
    element = {
        "id": f"img_{ObjectId()}",
        "type": "image",
        "content": image_url,
        "position": position or {"x": 0.5, "y": 0.1},
        "size": size or {"width": 0.4, "height": 0.6},
        "style": {},
        "alt_text": alt_text,
    }
    db = get_db()
    await db.slides.update_one(
        {"presentation_id": presentation_id, "index": slide_index},
        {"$push": {"elements": element}, "$set": {"updated_at": datetime.utcnow()}},
    )
    return {"element_id": element["id"], "image_url": image_url}


# ── QA Tools ──────────────────────────────────────────────────

async def snapshot(presentation_id: str) -> dict[str, Any]:
    """Take a snapshot of the current presentation state for diffing."""
    db = get_db()
    pres = await db.presentations.find_one({"_id": presentation_id})
    if not pres:
        raise ValueError(f"Presentation {presentation_id} not found")
    slides = await db.slides.find({"presentation_id": presentation_id}).sort("index", 1).to_list(200)
    return {
        "presentation": {k: v for k, v in pres.items() if k != "_id"},
        "slides": [{k: v for k, v in s.items() if k != "_id"} for s in slides],
        "timestamp": datetime.utcnow().isoformat(),
    }


async def screenshot(
    presentation_id: str,
    slide_index: int = 0,
) -> dict[str, Any]:
    """Capture a visual screenshot of a rendered slide (requires Playwright)."""
    return {
        "presentation_id": presentation_id,
        "slide_index": slide_index,
        "status": "screenshot_pending",
        "note": "Playwright render pipeline will capture in Phase 2",
    }


async def validate_layout(presentation_id: str) -> dict[str, Any]:
    """Validate that all slides have proper layout and content fit."""
    db = get_db()
    slides = await db.slides.find({"presentation_id": presentation_id}).sort("index", 1).to_list(200)
    issues: list[dict[str, Any]] = []
    for slide in slides:
        content = slide.get("content", {})
        title = content.get("title", "")
        if not title:
            issues.append({"slide_index": slide["index"], "issue": "Empty title"})
        bullets = content.get("bullets", [])
        if bullets and len(bullets) > 8:
            issues.append({
                "slide_index": slide["index"],
                "issue": f"Too many bullets ({len(bullets)}), recommend max 8",
            })
        elements = slide.get("elements", [])
        for elem in elements:
            pos = elem.get("position", {})
            if pos.get("x", 0) + elem.get("size", {}).get("width", 0) > 1.0:
                issues.append({
                    "slide_index": slide["index"],
                    "issue": f"Element '{elem.get('id')}' exceeds slide width",
                })
    return {"valid": len(issues) == 0, "issues": issues, "slides_checked": len(slides)}


async def score_quality(presentation_id: str) -> dict[str, Any]:
    """Score presentation quality across multiple dimensions."""
    strategy = await validate_strategy(presentation_id)
    layout = await validate_layout(presentation_id)

    scores = {
        "strategy": strategy["score"],
        "layout": 100 if layout["valid"] else max(0, 100 - len(layout["issues"]) * 10),
        "content_completeness": strategy["analysis"]["slide_count"] * 10 if strategy["analysis"]["slide_count"] <= 10 else 100,
    }
    overall = sum(scores.values()) // len(scores)
    return {
        "overall": overall,
        "scores": scores,
        "strategy_issues": strategy["issues"],
        "layout_issues": layout["issues"],
    }


# ── VFX Tools ─────────────────────────────────────────────────

async def create_3d_scene(
    slide_id: str,
    scene_type: str = "particles",
    config: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    """Configure a Three.js 3D scene for a slide."""
    db = get_db()
    scene = {
        "type": scene_type,
        "data": {},
        "config": config or {},
    }
    await db.slides.update_one(
        {"_id": slide_id},
        {"$set": {"threeScene": scene, "updated_at": datetime.utcnow()}},
    )
    return {"slide_id": slide_id, "scene_type": scene_type}


async def add_particles(
    slide_id: str,
    count: int = 100,
    color: str = "#ffffff",
    speed: float = 1.0,
) -> dict[str, Any]:
    """Add a particle system to a slide's 3D scene."""
    config = {"count": count, "color": color, "speed": speed}
    return await create_3d_scene(slide_id, scene_type="particles", config=config)


async def animate_element(
    slide_id: str,
    element_id: str,
    animation: str = "fade-in",
    order: int = 0,
    delay: int = 0,
) -> dict[str, Any]:
    """Add a fragment animation to a slide element."""
    db = get_db()
    fragment = {
        "elementId": element_id,
        "order": order,
        "animation": animation,
        "delay": delay,
    }
    await db.slides.update_one(
        {"_id": slide_id},
        {"$push": {"fragments": fragment}, "$set": {"updated_at": datetime.utcnow()}},
    )
    return {"slide_id": slide_id, "element_id": element_id, "animation": animation}


async def create_globe(
    slide_id: str,
    data_points: Optional[list[dict[str, Any]]] = None,
    config: Optional[dict[str, Any]] = None,
) -> dict[str, Any]:
    """Create a 3D globe visualization on a slide."""
    scene_config = {
        "data_points": data_points or [],
        **(config or {}),
    }
    return await create_3d_scene(slide_id, scene_type="globe", config=scene_config)


# ── Layout Tools ──────────────────────────────────────────────

async def select_layout(
    slide_type: str,
    content_length: int = 0,
    has_image: bool = False,
    has_chart: bool = False,
) -> dict[str, Any]:
    """Recommend an optimal layout for a slide based on content analysis."""
    if has_chart:
        layout = "chart"
    elif has_image and content_length > 200:
        layout = "text-left-visual-right"
    elif has_image:
        layout = "full-bleed"
    elif content_length > 500:
        layout = "bullets"
    elif slide_type == "title-slide":
        layout = "center-focus"
    elif slide_type == "comparison-slide" or slide_type == "competition-slide":
        layout = "comparison"
    elif slide_type == "team-slide":
        layout = "team-grid"
    else:
        layout = "center-focus"
    return {
        "recommended_layout": layout,
        "slide_type": slide_type,
        "reasoning": f"Selected {layout} for {slide_type} with {content_length} chars",
    }


# ═══════════════════════════════════════════════════════════════════
# V2 EVIDENCE-BASED CONTENT TOOLS
# ═══════════════════════════════════════════════════════════════════


class V2ResearchRequest(BaseModel):
    """Request for V2 evidence-based research."""
    topic: str = Field(..., min_length=1, max_length=500)
    slide_kind: str = Field(default="problem")
    audience: str = Field(default="investors", max_length=100)
    budget_mode: str = Field(default="lean", pattern=r"^(lean|balanced|hero)$")


class V2GenerateContentRequest(BaseModel):
    """Request for V2 slide content generation."""
    deck_id: str
    topic: str = Field(..., min_length=1, max_length=500)
    audience: str = Field(default="investors")
    budget_mode: str = Field(default="lean")
    style: str = Field(default="yc_crisp")
    outline: dict = Field(default_factory=dict)


async def plan_slide_research(
    topic: str,
    slide_title: str,
    slide_kind: str = "problem",
    audience: str = "investors",
) -> dict[str, Any]:
    """Plan research queries for a slide using QueryPlanner."""
    from app.mcp.brain_mcp.research.query_planner import QueryPlanner
    from app.mcp.brain_mcp.research.models import SlideKind

    try:
        kind = SlideKind(slide_kind)
    except ValueError:
        kind = SlideKind.PROBLEM

    planner = QueryPlanner()
    queries = await planner.plan_queries(topic, slide_title, kind, audience)
    return {
        "slide_kind": slide_kind,
        "queries": queries,
        "query_count": len(queries),
    }


async def run_slide_research(
    slide_id: str,
    topic: str,
    slide_kind: str = "problem",
    budget_mode: str = "lean",
) -> dict[str, Any]:
    """Execute research for a single slide using the full pipeline."""
    from app.config import settings
    from app.mcp.brain_mcp.research.models import SlideKind, BudgetMode
    from app.mcp.brain_mcp.research.provider_registry import ProviderRegistry
    from app.mcp.brain_mcp.research.circuit_breaker import CircuitBreaker
    from app.mcp.brain_mcp.research.content_events import ContentEventEmitter
    from app.mcp.brain_mcp.research.research_router import ResearchRouter
    from app.mcp.brain_mcp.research.query_planner import QueryPlanner
    from app.mcp.brain_mcp.research.evidence_assembler import EvidenceAssembler

    try:
        kind = SlideKind(slide_kind)
    except ValueError:
        kind = SlideKind.PROBLEM

    registry = ProviderRegistry(settings)
    breaker = CircuitBreaker(None)
    emitter = ContentEventEmitter(slide_id, None)
    router = ResearchRouter(registry, breaker, emitter)
    planner = QueryPlanner()
    assembler = EvidenceAssembler()

    queries = await planner.plan_queries(topic, "", kind, "investors")
    packets = await router.research_slide(
        slide_id=slide_id,
        slide_kind=kind,
        queries=queries,
        topic=topic,
        budget_mode=BudgetMode(budget_mode),
    )
    bundle = assembler.assemble(slide_id, kind, packets)
    return {
        "slide_id": slide_id,
        "fact_count": len(packets),
        "evidence_score": bundle.evidence_score,
        "missing_data": [m.__dict__ if hasattr(m, '__dict__') else str(m) for m in bundle.missing_data],
        "top_facts": [p.to_dict() for p in packets[:5]],
    }


async def run_pitch_debate(
    deck_id: str,
    topic: str,
    slide_kind: str = "market",
) -> dict[str, Any]:
    """Run a multi-agent debate on slide evidence."""
    from app.mcp.brain_mcp.research.models import SlideKind, BudgetMode
    from app.mcp.brain_mcp.research.debate_loop import DebateLoop
    from app.mcp.brain_mcp.research.evidence_assembler import EvidenceAssembler
    from app.mcp.brain_mcp.research.content_events import ContentEventEmitter
    from app.services.llm.model_router import ModelRouter

    try:
        kind = SlideKind(slide_kind)
    except ValueError:
        kind = SlideKind.MARKET

    model_router = ModelRouter.get_instance()
    emitter = ContentEventEmitter(deck_id, None)
    debate = DebateLoop(model_router, emitter)

    # Create a minimal bundle for debate
    assembler = EvidenceAssembler()
    bundle = assembler.assemble(f"debate_{slide_kind}", kind, [])

    outcome = await debate.run_debate(bundle, topic, kind)
    return {
        "confidence": outcome.confidence,
        "approved_claims": outcome.approved_claims,
        "rejected_claims": [r.__dict__ if hasattr(r, '__dict__') else str(r) for r in outcome.rejected_claims],
        "rounds": outcome.rounds_used,
    }


async def generate_slide_content_v2(
    deck_id: str,
    topic: str,
    slide_kind: str = "problem",
    audience: str = "investors",
    style: str = "yc_crisp",
    budget_mode: str = "lean",
) -> dict[str, Any]:
    """Generate evidence-based slide content using V2 pipeline."""
    from app.tasks.research_tasks import generate_deck_content

    outline = {
        "slides": [{"id": f"slide_0", "kind": slide_kind, "title": topic}]
    }

    task = generate_deck_content.delay(
        deck_id=deck_id,
        outline=outline,
        budget_mode=budget_mode,
        style=style,
        topic=topic,
        audience=audience,
    )
    return {
        "task_id": task.id,
        "status": "started",
        "deck_id": deck_id,
        "message": "Generation started in background. Use get-deck-run-status to check progress.",
    }


async def verify_slide_claims(
    claims: list[dict[str, Any]],
    evidence: list[dict[str, Any]],
) -> dict[str, Any]:
    """Verify claims against evidence using CitationGuard."""
    from app.mcp.brain_mcp.research.citation_guard import CitationGuard
    from app.mcp.brain_mcp.research.models import FactPacket

    guard = CitationGuard()
    packets = []
    for e in evidence:
        try:
            packets.append(FactPacket.from_dict(e))
        except Exception:
            continue

    results = []
    for claim in claims:
        text = claim.get("text", "")
        matched = guard.find_supporting_evidence(text, packets)
        results.append({
            "claim": text,
            "verified": len(matched) > 0,
            "supporting_sources": len(matched),
        })
    return {"verified_claims": results, "total": len(results)}


async def get_evidence_graph(deck_id: str) -> dict[str, Any]:
    """Get the evidence graph for a deck run."""
    db = get_db()
    run = await db.deck_runs.find_one(
        {"deck_id": deck_id},
        {"evidence_graph": 1}
    )
    if not run:
        return {"error": "Deck run not found", "graph": None}
    return {"graph": run.get("evidence_graph", {})}


async def get_deck_run_status(deck_id: str) -> dict[str, Any]:
    """Get the generation status for a deck run."""
    db = get_db()
    run = await db.deck_runs.find_one(
        {"deck_id": deck_id},
        {"status": 1, "errors": 1, "total_time_ms": 1, "total_fact_packets": 1}
    )
    if not run:
        return {"status": "not_found", "deck_id": deck_id}
    run.pop("_id", None)
    return run


async def optimize_grid(
    element_count: int,
    aspect_ratio: str = "16:9",
) -> dict[str, Any]:
    """Calculate optimal grid layout for N elements."""
    if element_count <= 2:
        cols, rows = element_count, 1
    elif element_count <= 4:
        cols, rows = 2, 2
    elif element_count <= 6:
        cols, rows = 3, 2
    elif element_count <= 9:
        cols, rows = 3, 3
    else:
        cols, rows = 4, (element_count + 3) // 4
    return {
        "cols": cols,
        "rows": rows,
        "element_count": element_count,
        "aspect_ratio": aspect_ratio,
    }


async def measure_content(content: dict[str, Any]) -> dict[str, Any]:
    """Estimate content dimensions for layout optimisation."""
    title_len = len(content.get("title", ""))
    body_len = len(content.get("body_text", "") or "")
    bullet_count = len(content.get("bullets", []) or [])
    has_image = bool(content.get("image_url"))
    has_chart = bool(content.get("chart_data"))

    estimated_height = 0.1  # title zone
    if title_len > 0:
        estimated_height += 0.1
    if body_len > 0:
        estimated_height += min(0.5, body_len / 500.0)
    if bullet_count > 0:
        estimated_height += min(0.4, bullet_count * 0.05)

    return {
        "title_length": title_len,
        "body_length": body_len,
        "bullet_count": bullet_count,
        "has_image": has_image,
        "has_chart": has_chart,
        "estimated_height": round(estimated_height, 2),
        "overflow_risk": estimated_height > 0.85,
    }


async def validate_fit(
    slide_id: str,
) -> dict[str, Any]:
    """Validate that slide content fits within the layout constraints."""
    db = get_db()
    slide = await db.slides.find_one({"_id": slide_id})
    if not slide:
        raise ValueError(f"Slide {slide_id} not found")

    content = slide.get("content", {})
    measurement = await measure_content(content)
    elements = slide.get("elements", [])

    overlap_pairs: list[tuple[str, str]] = []
    for i, a in enumerate(elements):
        for b in elements[i + 1:]:
            a_pos = a.get("position", {})
            a_size = a.get("size", {})
            b_pos = b.get("position", {})
            b_size = b.get("size", {})
            if (
                a_pos.get("x", 0) < b_pos.get("x", 0) + b_size.get("width", 0)
                and a_pos.get("x", 0) + a_size.get("width", 0) > b_pos.get("x", 0)
                and a_pos.get("y", 0) < b_pos.get("y", 0) + b_size.get("height", 0)
                and a_pos.get("y", 0) + a_size.get("height", 0) > b_pos.get("y", 0)
            ):
                overlap_pairs.append((a.get("id", "?"), b.get("id", "?")))

    return {
        "slide_id": slide_id,
        "fits": not measurement["overflow_risk"] and len(overlap_pairs) == 0,
        "overflow_risk": measurement["overflow_risk"],
        "overlapping_elements": overlap_pairs,
    }


# ═══════════════════════════════════════════════════════════════════
# V3 UNIFIED PIPELINE TOOLS
# ═══════════════════════════════════════════════════════════════════


class V3GenerateToolRequest(BaseModel):
    topic: str = Field(..., min_length=1, max_length=500)
    description: str = Field(default="", max_length=5000)
    audience: str = Field(default="investors", max_length=200)
    purpose: str = Field(default="pitch", max_length=100)
    mode: str = Field(default="standard", pattern=r"^(standard|premium)$")
    slide_count: int = Field(default=10, ge=3, le=30)
    writing_style: str = Field(default="yc_crisp", max_length=50)


async def generate_deck_v3(
    topic: str,
    description: str = "",
    audience: str = "investors",
    purpose: str = "pitch",
    mode: str = "standard",
    slide_count: int = 10,
    writing_style: str = "yc_crisp",
) -> dict[str, Any]:
    """Start a V3 unified deck generation (standard or premium mode)."""
    import uuid
    from app.tasks.unified_tasks import generate_unified_deck

    deck_id = str(uuid.uuid4())
    queue = "content-fast" if mode == "standard" else "content-premium"
    request_dict = {
        "topic": topic,
        "description": description,
        "audience": audience,
        "purpose": purpose,
        "mode": mode,
        "slide_count": slide_count,
        "writing_style": writing_style,
        "language": "en",
        "generate_notes": mode == "premium",
        "target_formats": ["revealjs"],
    }
    task = generate_unified_deck.apply_async(
        args=[deck_id, request_dict], queue=queue
    )
    return {"deck_id": deck_id, "task_id": task.id, "mode": mode, "status": "queued"}


async def get_deck_v3_status(deck_id: str) -> dict[str, Any]:
    """Check the status of a V3 deck generation."""
    db = get_db()
    run = await db.deck_runs_v3.find_one({"deck_id": deck_id})
    if not run:
        return {"error": f"Deck {deck_id} not found"}
    run.pop("_id", None)
    return {
        "deck_id": deck_id,
        "status": run.get("status", "unknown"),
        "mode": run.get("mode", "standard"),
        "total_slides": len(run.get("slides", [])),
        "quality_score": run.get("quality_score", 0.0),
        "total_time_ms": run.get("total_time_ms", 0.0),
    }


async def get_deck_v3_evidence(deck_id: str) -> dict[str, Any]:
    """Get the evidence report for a premium V3 deck."""
    db = get_db()
    run = await db.deck_runs_v3.find_one({"deck_id": deck_id})
    if not run:
        return {"error": f"Deck {deck_id} not found"}
    if run.get("mode") != "premium":
        return {"error": "Evidence only available for premium mode"}
    return {
        "deck_id": deck_id,
        "evidence_report": run.get("evidence_report", {}),
        "coherence_score": run.get("coherence_score", 0.0),
    }


async def cancel_deck_v3(deck_id: str) -> dict[str, Any]:
    """Cancel an in-progress V3 deck generation."""
    import redis
    from app.config import settings

    try:
        r = redis.from_url(settings.REDIS_URL)
        r.set(f"deck:{deck_id}:cancel", "1", ex=600)
        r.close()
    except Exception as e:
        logger.warning("cancel_redis_failed", error=str(e))

    db = get_db()
    await db.deck_runs_v3.update_one(
        {"deck_id": deck_id, "status": "running"},
        {"$set": {"status": "cancelling"}},
    )
    return {"deck_id": deck_id, "status": "cancelling"}


# ═══════════════════════════════════════════════════════════════════
# REGISTRY — maps tool names to handlers for FastMCP registration
# ═══════════════════════════════════════════════════════════════════

TOOL_REGISTRY: dict[str, dict[str, Any]] = {
    # Presentation Management
    "create_presentation": {"handler": create_presentation, "category": "presentation"},
    "get_presentation": {"handler": get_presentation, "category": "presentation"},
    "update_presentation": {"handler": update_presentation, "category": "presentation"},
    "delete_presentation": {"handler": delete_presentation, "category": "presentation"},
    "list_presentations": {"handler": list_presentations, "category": "presentation"},
    "duplicate_presentation": {"handler": duplicate_presentation, "category": "presentation"},
    # Slide Operations
    "add_slide": {"handler": add_slide, "category": "slide"},
    "update_slide": {"handler": update_slide, "category": "slide"},
    "delete_slide": {"handler": delete_slide, "category": "slide"},
    "reorder_slides": {"handler": reorder_slides, "category": "slide"},
    "get_slide": {"handler": get_slide, "category": "slide"},
    "get_all_slides": {"handler": get_all_slides, "category": "slide"},
    "duplicate_slide": {"handler": duplicate_slide, "category": "slide"},
    # Content Generation
    "generate_dsl": {"handler": generate_dsl, "category": "content"},
    "compile_react": {"handler": compile_react, "category": "content"},
    "compile_revealjs": {"handler": compile_revealjs, "category": "content"},
    "compile_html": {"handler": compile_html, "category": "content"},
    "generate_outline": {"handler": generate_outline, "category": "content"},
    "refine_content": {"handler": refine_content, "category": "content"},
    # Strategy
    "analyze_presentation": {"handler": analyze_presentation, "category": "strategy"},
    "validate_strategy": {"handler": validate_strategy, "category": "strategy"},
    "select_archetype": {"handler": select_archetype, "category": "strategy"},
    # Research
    "research_topic": {"handler": research_topic, "category": "research"},
    "extract_document": {"handler": extract_document, "category": "research"},
    "search_web": {"handler": search_web, "category": "research"},
    "analyze_data": {"handler": analyze_data, "category": "research"},
    # Design
    "apply_theme": {"handler": apply_theme, "category": "design"},
    "generate_theme": {"handler": generate_theme, "category": "design"},
    "discover_style": {"handler": discover_style, "category": "design"},
    "check_contrast": {"handler": check_contrast, "category": "design"},
    # PPTX
    "create_pptx": {"handler": create_pptx, "category": "pptx"},
    "add_chart": {"handler": add_chart, "category": "pptx"},
    "add_table": {"handler": add_table, "category": "pptx"},
    "add_shape": {"handler": add_shape, "category": "pptx"},
    "add_image": {"handler": add_image, "category": "pptx"},
    # QA
    "snapshot": {"handler": snapshot, "category": "qa"},
    "screenshot": {"handler": screenshot, "category": "qa"},
    "validate_layout": {"handler": validate_layout, "category": "qa"},
    "score_quality": {"handler": score_quality, "category": "qa"},
    # VFX
    "create_3d_scene": {"handler": create_3d_scene, "category": "vfx"},
    "add_particles": {"handler": add_particles, "category": "vfx"},
    "animate_element": {"handler": animate_element, "category": "vfx"},
    "create_globe": {"handler": create_globe, "category": "vfx"},
    # Layout
    "select_layout": {"handler": select_layout, "category": "layout"},
    "optimize_grid": {"handler": optimize_grid, "category": "layout"},
    "measure_content": {"handler": measure_content, "category": "layout"},
    "validate_fit": {"handler": validate_fit, "category": "layout"},
    # V3 Unified Pipeline
    "generate_deck_v3": {"handler": generate_deck_v3, "category": "v3-unified"},
    "get_deck_v3_status": {"handler": get_deck_v3_status, "category": "v3-unified"},
    "get_deck_v3_evidence": {"handler": get_deck_v3_evidence, "category": "v3-unified"},
    "cancel_deck_v3": {"handler": cancel_deck_v3, "category": "v3-unified"},
}


def get_tool_names_by_category(category: str) -> list[str]:
    """Return tool names belonging to a category."""
    return [
        name for name, meta in TOOL_REGISTRY.items()
        if meta["category"] == category
    ]


def get_all_categories() -> list[str]:
    """Return unique category names."""
    return sorted({meta["category"] for meta in TOOL_REGISTRY.values()})
