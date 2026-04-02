"""
PPTX Builder — generates PowerPoint files with native Excel-backed charts.
Uses python-pptx for full PPTX creation with editable charts.
"""

import io
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

import structlog

from app.mcp.render_mcp.config import (
    PPTX_SLIDE_WIDTH_EMU,
    PPTX_SLIDE_HEIGHT_EMU,
    PPTX_MARGIN_EMU,
)

logger = structlog.get_logger()


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class PptxBuilder:
    """Builds PPTX presentations with native charts and themed styling."""

    CHART_TYPE_MAP = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "donut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
    }

    def build(self, slides: list[dict], theme: dict, metadata: dict = None) -> bytes:
        """Build a complete PPTX file. Returns bytes."""
        prs = Presentation()
        prs.slide_width = Emu(PPTX_SLIDE_WIDTH_EMU)
        prs.slide_height = Emu(PPTX_SLIDE_HEIGHT_EMU)

        colors = theme.get("colors", {})
        fonts = theme.get("fonts", {})
        heading_font = fonts.get("heading", "Calibri")
        body_font = fonts.get("body", "Calibri")

        for i, slide_data in enumerate(slides):
            layout = slide_data.get("layout", "bullets")
            content = slide_data.get("content", {})

            pptx_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

            # Set slide background
            self._set_background(pptx_slide, colors, layout)

            # Build per layout type
            builder_method = getattr(self, f"_build_{layout.replace('-', '_')}", None)
            if builder_method:
                builder_method(pptx_slide, content, colors, heading_font, body_font)
            else:
                self._build_bullets(pptx_slide, content, colors, heading_font, body_font)

            # Add speaker notes
            notes = slide_data.get("speaker_notes", "")
            if notes:
                notes_slide = pptx_slide.notes_slide
                notes_slide.notes_text_frame.text = notes

        # Save to bytes
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        logger.info("pptx_built", slide_count=len(slides))
        return buf.read()

    def _set_background(self, slide, colors: dict, layout: str):
        """Set slide background color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        bg_color = colors.get("background", "#ffffff")
        if layout == "title-hero":
            bg_color = colors.get("primary", bg_color)
        fill.fore_color.rgb = _hex_to_rgb(bg_color)

    def _add_title(self, slide, text: str, colors: dict, font_name: str,
                   left=None, top=None, width=None, height=None, font_size=36,
                   is_hero=False):
        """Add a title text box."""
        left = left or Emu(PPTX_MARGIN_EMU)
        top = top or Emu(PPTX_MARGIN_EMU)
        width = width or Emu(PPTX_SLIDE_WIDTH_EMU - 2 * PPTX_MARGIN_EMU)
        height = height or Inches(1.2)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.name = font_name
        if is_hero:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        else:
            p.font.color.rgb = _hex_to_rgb(colors.get("text_primary", "#111827"))
        return txBox

    def _add_body_text(self, slide, text: str, colors: dict, font_name: str,
                       left=None, top=None, width=None, height=None, font_size=16):
        """Add body text box."""
        left = left or Emu(PPTX_MARGIN_EMU)
        top = top or Inches(2)
        width = width or Emu(PPTX_SLIDE_WIDTH_EMU - 2 * PPTX_MARGIN_EMU)
        height = height or Inches(4)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = _hex_to_rgb(colors.get("text_secondary", "#4b5563"))
        return txBox

    # ── Layout builders ─────────────────────────────────────────

    def _build_title_hero(self, slide, content, colors, heading_font, body_font):
        self._add_title(
            slide, content.get("title", ""),
            colors, heading_font,
            top=Inches(2.5), font_size=44, is_hero=True,
        )
        subtitle = content.get("subtitle", "")
        if subtitle:
            txBox = slide.shapes.add_textbox(
                Emu(PPTX_MARGIN_EMU), Inches(4),
                Emu(PPTX_SLIDE_WIDTH_EMU - 2 * PPTX_MARGIN_EMU), Inches(1),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.name = body_font
            p.font.color.rgb = RGBColor(220, 220, 220)
            p.alignment = PP_ALIGN.CENTER

    def _build_bullets(self, slide, content, colors, heading_font, body_font):
        self._add_title(slide, content.get("title", ""), colors, heading_font, font_size=32)
        bullets = content.get("bullets", [])
        if bullets:
            txBox = slide.shapes.add_textbox(
                Emu(PPTX_MARGIN_EMU), Inches(1.8),
                Emu(PPTX_SLIDE_WIDTH_EMU - 2 * PPTX_MARGIN_EMU), Inches(5),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(18)
                p.font.name = body_font
                p.font.color.rgb = _hex_to_rgb(colors.get("text_primary", "#111827"))
                p.space_after = Pt(8)

    def _build_bullets_with_image(self, slide, content, colors, heading_font, body_font):
        self._add_title(slide, content.get("title", ""), colors, heading_font, font_size=32)
        # Left side bullets
        bullets = content.get("bullets", [])
        if bullets:
            half_width = Emu((PPTX_SLIDE_WIDTH_EMU - 3 * PPTX_MARGIN_EMU) // 2)
            txBox = slide.shapes.add_textbox(
                Emu(PPTX_MARGIN_EMU), Inches(1.8), half_width, Inches(5),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(16)
                p.font.name = body_font
                p.font.color.rgb = _hex_to_rgb(colors.get("text_primary", "#111827"))
                p.space_after = Pt(6)
        # Right side: image placeholder
        right_left = Emu(PPTX_SLIDE_WIDTH_EMU // 2 + PPTX_MARGIN_EMU // 2)
        placeholder = slide.shapes.add_textbox(
            right_left, Inches(1.8),
            Emu(PPTX_SLIDE_WIDTH_EMU // 2 - PPTX_MARGIN_EMU * 2), Inches(4.5),
        )
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Image: {content.get('image_prompt', 'Image placeholder')}]"
        p.font.size = Pt(12)
        p.font.color.rgb = _hex_to_rgb(colors.get("text_secondary", "#9ca3af"))
        p.alignment = PP_ALIGN.CENTER

    def _build_two_column(self, slide, content, colors, heading_font, body_font):
        self._add_title(slide, content.get("title", ""), colors, heading_font, font_size=32)
        half_width = Emu((PPTX_SLIDE_WIDTH_EMU - 3 * PPTX_MARGIN_EMU) // 2)

        for col, key, x_offset in [
            ("left", "left_content", Emu(PPTX_MARGIN_EMU)),
            ("right", "right_content", Emu(PPTX_SLIDE_WIDTH_EMU // 2 + PPTX_MARGIN_EMU // 2)),
        ]:
            text = content.get(key, "")
            if text:
                self._add_body_text(
                    slide, text, colors, body_font,
                    left=x_offset, top=Inches(1.8),
                    width=half_width, height=Inches(5),
                )

    def _build_chart(self, slide, content, colors, heading_font, body_font):
        self._add_title(slide, content.get("title", ""), colors, heading_font, font_size=28)

        chart_data_raw = content.get("chart_data", {})
        chart_type_str = content.get("chart_type", "bar")
        xl_chart_type = self.CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        labels = chart_data_raw.get("labels", ["A", "B", "C"])
        datasets = chart_data_raw.get("datasets", [{"label": "Data", "values": [1, 2, 3]}])

        chart_data = CategoryChartData()
        chart_data.categories = labels
        for ds in datasets:
            chart_data.add_series(ds.get("label", "Series"), ds.get("values", []))

        chart_frame = slide.shapes.add_chart(
            xl_chart_type,
            Emu(PPTX_MARGIN_EMU), Inches(1.6),
            Emu(PPTX_SLIDE_WIDTH_EMU - 2 * PPTX_MARGIN_EMU), Inches(5),
            chart_data,
        )

        # Style the chart with theme colors
        chart = chart_frame.chart
        chart.has_legend = len(datasets) > 1
        chart_colors_list = colors.get("chart_colors", ["#2563eb", "#7c3aed", "#059669"])

        for idx, series in enumerate(chart.series):
            color = chart_colors_list[idx % len(chart_colors_list)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _hex_to_rgb(color)

        # Source attribution
        source = content.get("source_attribution", "")
        if source:
            txBox = slide.shapes.add_textbox(
                Emu(PPTX_MARGIN_EMU), Inches(6.8),
                Emu(PPTX_SLIDE_WIDTH_EMU - 2 * PPTX_MARGIN_EMU), Inches(0.3),
            )
            p = txBox.text_frame.paragraphs[0]
            p.text = f"Source: {source}"
            p.font.size = Pt(9)
            p.font.color.rgb = _hex_to_rgb(colors.get("text_secondary", "#9ca3af"))

    def _build_comparison(self, slide, content, colors, heading_font, body_font):
        self._add_title(slide, content.get("title", ""), colors, heading_font, font_size=32)
        half_width = Emu((PPTX_SLIDE_WIDTH_EMU - 3 * PPTX_MARGIN_EMU) // 2)

        for label_key, items_key, x_offset in [
            ("left_label", "left_items", Emu(PPTX_MARGIN_EMU)),
            ("right_label", "right_items", Emu(PPTX_SLIDE_WIDTH_EMU // 2 + PPTX_MARGIN_EMU // 2)),
        ]:
            label = content.get(label_key, "")
            items = content.get(items_key, [])

            # Label
            lbl_box = slide.shapes.add_textbox(x_offset, Inches(1.6), half_width, Inches(0.5))
            p = lbl_box.text_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.name = heading_font
            p.font.color.rgb = _hex_to_rgb(colors.get("primary", "#2563eb"))

            # Items
            if items:
                txBox = slide.shapes.add_textbox(x_offset, Inches(2.2), half_width, Inches(4.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                for j, item in enumerate(items):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = f"✓ {item}"
                    p.font.size = Pt(16)
                    p.font.name = body_font
                    p.font.color.rgb = _hex_to_rgb(colors.get("text_primary", "#111827"))
                    p.space_after = Pt(6)

    def _build_timeline(self, slide, content, colors, heading_font, body_font):
        self._add_title(slide, content.get("title", ""), colors, heading_font, font_size=32)
        events = content.get("events", [])
        if not events:
            return

        usable_width = PPTX_SLIDE_WIDTH_EMU - 2 * PPTX_MARGIN_EMU
        event_width = usable_width // max(len(events), 1)

        for i, event in enumerate(events):
            x = Emu(PPTX_MARGIN_EMU + i * event_width)
            # Date
            date_box = slide.shapes.add_textbox(x, Inches(2.2), Emu(event_width), Inches(0.5))
            p = date_box.text_frame.paragraphs[0]
            p.text = event.get("date", "")
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.name = heading_font
            p.font.color.rgb = _hex_to_rgb(colors.get("primary", "#2563eb"))
            p.alignment = PP_ALIGN.CENTER

            # Description
            desc_box = slide.shapes.add_textbox(x, Inches(2.8), Emu(event_width), Inches(3))
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = event.get("description", "")
            p.font.size = Pt(13)
            p.font.name = body_font
            p.font.color.rgb = _hex_to_rgb(colors.get("text_primary", "#111827"))
            p.alignment = PP_ALIGN.CENTER

    def _build_quote(self, slide, content, colors, heading_font, body_font):
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _hex_to_rgb(colors.get("surface", "#f9fafb"))

        # Quote text
        txBox = slide.shapes.add_textbox(
            Inches(1.5), Inches(2),
            Emu(PPTX_SLIDE_WIDTH_EMU - Inches(3).emu), Inches(3),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{content.get("quote_text", "")}"'
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.name = heading_font
        p.font.color.rgb = _hex_to_rgb(colors.get("text_primary", "#111827"))
        p.alignment = PP_ALIGN.CENTER

        # Author
        author = content.get("quote_author", "")
        role = content.get("quote_role", "")
        if author:
            auth_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(5),
                Emu(PPTX_SLIDE_WIDTH_EMU - Inches(3).emu), Inches(0.8),
            )
            p = auth_box.text_frame.paragraphs[0]
            p.text = f"— {author}" + (f", {role}" if role else "")
            p.font.size = Pt(16)
            p.font.name = body_font
            p.font.color.rgb = _hex_to_rgb(colors.get("primary", "#2563eb"))
            p.alignment = PP_ALIGN.CENTER

    def _build_team_grid(self, slide, content, colors, heading_font, body_font):
        self._add_title(slide, content.get("title", ""), colors, heading_font, font_size=32)
        members = content.get("members", [])
        if not members:
            return

        cols = min(len(members), 3)
        rows = (len(members) + cols - 1) // cols
        card_width = (PPTX_SLIDE_WIDTH_EMU - (cols + 1) * PPTX_MARGIN_EMU) // cols
        card_height = Inches(1.8)

        for idx, member in enumerate(members):
            col = idx % cols
            row = idx // cols
            x = Emu(PPTX_MARGIN_EMU + col * (card_width + PPTX_MARGIN_EMU))
            y = Emu(Inches(1.8).emu + row * (card_height.emu + Inches(0.3).emu))

            txBox = slide.shapes.add_textbox(x, y, Emu(card_width), card_height)
            tf = txBox.text_frame
            tf.word_wrap = True

            # Name
            p = tf.paragraphs[0]
            p.text = member.get("name", "")
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.name = heading_font
            p.font.color.rgb = _hex_to_rgb(colors.get("text_primary", "#111827"))

            # Role
            p2 = tf.add_paragraph()
            p2.text = member.get("role", "")
            p2.font.size = Pt(13)
            p2.font.name = body_font
            p2.font.color.rgb = _hex_to_rgb(colors.get("primary", "#2563eb"))

            # Bio
            bio = member.get("bio", "")
            if bio:
                p3 = tf.add_paragraph()
                p3.text = bio
                p3.font.size = Pt(11)
                p3.font.name = body_font
                p3.font.color.rgb = _hex_to_rgb(colors.get("text_secondary", "#6b7280"))

    def _build_kpi_dashboard(self, slide, content, colors, heading_font, body_font):
        self._add_title(slide, content.get("title", ""), colors, heading_font, font_size=32)
        metrics = content.get("metrics", [])
        if not metrics:
            return

        cols = min(len(metrics), 3)
        card_width = (PPTX_SLIDE_WIDTH_EMU - (cols + 1) * PPTX_MARGIN_EMU) // cols
        rows = (len(metrics) + cols - 1) // cols

        for idx, metric in enumerate(metrics):
            col = idx % cols
            row = idx // cols
            x = Emu(PPTX_MARGIN_EMU + col * (card_width + PPTX_MARGIN_EMU))
            y = Emu(Inches(1.8).emu + row * Inches(2.2).emu)

            txBox = slide.shapes.add_textbox(x, y, Emu(card_width), Inches(1.8))
            tf = txBox.text_frame
            tf.word_wrap = True

            # Value (large)
            p = tf.paragraphs[0]
            p.text = str(metric.get("value", ""))
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.name = heading_font
            p.font.color.rgb = _hex_to_rgb(colors.get("primary", "#2563eb"))
            p.alignment = PP_ALIGN.CENTER

            # Change indicator
            change = metric.get("change", "")
            if change:
                p2 = tf.add_paragraph()
                p2.text = change
                p2.font.size = Pt(16)
                p2.font.name = body_font
                is_positive = change.startswith("+") or "up" in change.lower()
                p2.font.color.rgb = _hex_to_rgb("#059669" if is_positive else "#dc2626")
                p2.alignment = PP_ALIGN.CENTER

            # Label
            p3 = tf.add_paragraph()
            p3.text = metric.get("label", "")
            p3.font.size = Pt(13)
            p3.font.name = body_font
            p3.font.color.rgb = _hex_to_rgb(colors.get("text_secondary", "#6b7280"))
            p3.alignment = PP_ALIGN.CENTER

    def _build_full_image(self, slide, content, colors, heading_font, body_font):
        # Placeholder for full-image layout (image added separately)
        title = content.get("title", "")
        subtitle = content.get("subtitle", "")

        # Image placeholder
        ph = slide.shapes.add_textbox(
            Inches(2), Inches(2),
            Emu(PPTX_SLIDE_WIDTH_EMU - Inches(4).emu), Inches(3),
        )
        p = ph.text_frame.paragraphs[0]
        p.text = f"[Image: {content.get('image_prompt', 'Full-bleed image')}]"
        p.font.size = Pt(14)
        p.font.color.rgb = _hex_to_rgb(colors.get("text_secondary", "#9ca3af"))
        p.alignment = PP_ALIGN.CENTER

        if title:
            self._add_title(slide, title, colors, heading_font, top=Inches(5.5), font_size=28)

    def _build_blank(self, slide, content, colors, heading_font, body_font):
        title = content.get("title", "")
        if title:
            self._add_title(slide, title, colors, heading_font, font_size=32)
        body = content.get("body_text", "")
        if body:
            self._add_body_text(slide, body, colors, body_font)
