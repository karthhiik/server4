"""
Chart Builder — creates native Excel-backed charts for PPTX
and Chart.js configs for HTML rendering.
"""

from typing import Optional

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

import structlog

logger = structlog.get_logger()


CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class ChartBuilder:
    """Creates charts for PPTX and HTML export."""

    def add_chart_to_slide(
        self,
        pptx_slide,
        chart_data: dict,
        chart_colors: list[str],
        left=None,
        top=None,
        width=None,
        height=None,
    ):
        """Add a native Excel-backed chart to a PPTX slide."""
        chart_type_str = chart_data.get("chart_type", "bar")
        xl_type = CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

        labels = chart_data.get("labels", [])
        datasets = chart_data.get("datasets", [])

        cd = CategoryChartData()
        cd.categories = labels
        for ds in datasets:
            cd.add_series(ds.get("label", "Series"), ds.get("values", []))

        left = left or Inches(0.5)
        top = top or Inches(1.5)
        width = width or Inches(12)
        height = height or Inches(5.5)

        chart_frame = pptx_slide.shapes.add_chart(xl_type, left, top, width, height, cd)
        chart = chart_frame.chart

        # Style with theme colors
        chart.has_legend = len(datasets) > 1
        if chart.has_legend:
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)

        for idx, series in enumerate(chart.series):
            color = chart_colors[idx % len(chart_colors)] if chart_colors else "#2563eb"
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _hex_to_rgb(color)

            # For pie/donut, color individual points
            if chart_type_str in ("pie", "donut"):
                for pt_idx in range(len(labels)):
                    pt_color = chart_colors[pt_idx % len(chart_colors)]
                    pt = series.points[pt_idx]
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = _hex_to_rgb(pt_color)

        return chart_frame

    def build_chartjs_config(self, chart_data: dict, chart_colors: list[str]) -> dict:
        """Build a Chart.js configuration object for HTML rendering."""
        chart_type = chart_data.get("chart_type", "bar")
        labels = chart_data.get("labels", [])
        datasets = chart_data.get("datasets", [])

        # Map our types to Chart.js types
        chartjs_type = {
            "bar": "bar",
            "stacked_bar": "bar",
            "line": "line",
            "pie": "pie",
            "donut": "doughnut",
            "area": "line",
        }.get(chart_type, "bar")

        chartjs_datasets = []
        for i, ds in enumerate(datasets):
            color = chart_colors[i % len(chart_colors)] if chart_colors else "#2563eb"
            cfg = {
                "label": ds.get("label", f"Series {i+1}"),
                "data": ds.get("values", []),
                "backgroundColor": color,
                "borderColor": color,
            }
            if chart_type == "area":
                cfg["fill"] = True
                cfg["backgroundColor"] = color + "26"  # 15% opacity
            if chart_type in ("pie", "donut"):
                cfg["backgroundColor"] = [
                    chart_colors[j % len(chart_colors)] for j in range(len(labels))
                ]
            chartjs_datasets.append(cfg)

        config = {
            "type": chartjs_type,
            "data": {"labels": labels, "datasets": chartjs_datasets},
            "options": {
                "responsive": True,
                "plugins": {"legend": {"display": len(datasets) > 1}},
            },
        }

        if chart_type == "stacked_bar":
            config["options"]["scales"] = {
                "x": {"stacked": True},
                "y": {"stacked": True},
            }

        return config
