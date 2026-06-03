"""
V4 screenshot-based exporters.

After ``slide_screenshot.capture_deck_screenshots`` produces one PNG per
slide, these helpers wrap the bytes into PDF, PPTX, and DOCX files. Each
format renders the screenshot as the only content of a page/slide so the
download matches exactly what the user reviewed in /studio.

Public API::

    pdf_bytes  = build_pdf_from_screenshots(pngs, title)
    pptx_bytes = build_pptx_from_screenshots(pngs, title)
    docx_bytes = build_docx_from_screenshots(pngs, title)

All three are pure-CPU once the screenshots are in hand, so they can be
called inline without the chromium worker thread.
"""

from __future__ import annotations

import io
from typing import Sequence

import structlog

logger = structlog.get_logger(__name__)


# ── PDF ──────────────────────────────────────────────────────────


def build_pdf_from_screenshots(pages: Sequence[bytes], title: str = "Presentation") -> bytes:
    """Assemble PNG screenshots into a 16:9 landscape PDF.

    Uses Pillow's PDF save mode which embeds PNGs losslessly. Each page
    matches the screenshot's pixel dimensions so no scaling distortion.
    Falls back to a minimal byte stream if Pillow is unavailable.
    """
    if not pages:
        raise ValueError("build_pdf_from_screenshots requires at least one page")

    try:
        from PIL import Image
    except ImportError:
        logger.warning("pillow_not_available_for_pdf")
        return b""

    images: list[Image.Image] = []
    for idx, png in enumerate(pages):
        try:
            img = Image.open(io.BytesIO(png)).convert("RGB")
            images.append(img)
        except Exception as exc:  # noqa: BLE001
            logger.warning("pdf_page_decode_failed", index=idx, error=str(exc)[:200])

    if not images:
        raise ValueError("All screenshot pages failed to decode")

    buf = io.BytesIO()
    images[0].save(
        buf,
        format="PDF",
        save_all=True,
        append_images=images[1:],
        title=title[:120],
        resolution=144.0,
    )
    return buf.getvalue()


# ── PPTX ─────────────────────────────────────────────────────────


def build_pptx_from_screenshots(pages: Sequence[bytes], title: str = "Presentation") -> bytes:
    """Wrap each screenshot as a full-bleed slide in a 16:9 deck.

    The frame is set to 13.333" x 7.5" (the standard Microsoft 16:9 size)
    so the file opens cleanly in PowerPoint, Keynote, and Google Slides.
    """
    if not pages:
        raise ValueError("build_pptx_from_screenshots requires at least one page")

    try:
        from pptx import Presentation
        from pptx.util import Emu, Inches
    except ImportError:
        logger.warning("python_pptx_not_available")
        return b""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Apply title to PPTX core properties so the file shows up correctly
    # in PowerPoint's "recent files" list.
    cp = prs.core_properties
    cp.title = title[:120]

    blank_layout = prs.slide_layouts[6]  # 6 = blank layout in default master

    for idx, png in enumerate(pages):
        slide = prs.slides.add_slide(blank_layout)
        try:
            img_stream = io.BytesIO(png)
            slide.shapes.add_picture(
                img_stream,
                left=Emu(0),
                top=Emu(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )
        except Exception as exc:  # noqa: BLE001
            logger.warning("pptx_page_embed_failed", index=idx, error=str(exc)[:200])

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ── DOCX ─────────────────────────────────────────────────────────


def build_docx_from_screenshots(pages: Sequence[bytes], title: str = "Presentation") -> bytes:
    """Build a DOCX where each page is the slide screenshot at full width.

    Uses python-docx. Each slide image gets its own page (page break
    between slides) with a small header showing the slide number.
    """
    if not pages:
        raise ValueError("build_docx_from_screenshots requires at least one page")

    try:
        from docx import Document
        from docx.enum.text import WD_BREAK
        from docx.shared import Inches
    except ImportError:
        logger.warning("python_docx_not_available")
        return b""

    doc = Document()

    # Set landscape for the whole doc so the 16:9 screenshots fit naturally.
    section = doc.sections[0]
    new_w, new_h = section.page_height, section.page_width
    section.page_width = new_w
    section.page_height = new_h
    section.left_margin = Inches(0.5)
    section.right_margin = Inches(0.5)
    section.top_margin = Inches(0.5)
    section.bottom_margin = Inches(0.5)

    title_para = doc.add_paragraph()
    title_run = title_para.add_run(title[:120])
    title_run.bold = True
    title_run.font.size = doc.styles["Heading 1"].font.size

    for idx, png in enumerate(pages):
        if idx > 0:
            page_break = doc.add_paragraph()
            page_break.add_run().add_break(WD_BREAK.PAGE)

        header = doc.add_paragraph()
        header_run = header.add_run(f"Slide {idx + 1}")
        header_run.bold = True

        try:
            stream = io.BytesIO(png)
            doc.add_picture(stream, width=Inches(9.5))
        except Exception as exc:  # noqa: BLE001
            logger.warning("docx_page_embed_failed", index=idx, error=str(exc)[:200])
            doc.add_paragraph(f"[slide {idx + 1} could not be rendered]")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()
