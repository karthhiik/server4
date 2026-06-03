"""
V4 Document Extractor — extract clean text from user-uploaded files.

Supported types (extension-routed):
  .pdf            → pypdf
  .docx           → python-docx
  .pptx           → python-pptx
  .txt / .md / .rst / .csv / .json → utf-8 read
  .url            → web scraping (firecrawl/tavily)

Returns an `ExtractedDocument` with cleaned page-level text and a single
concatenated body. The pipeline can then chunk + embed via `embeddings.py`.

Why pypdf over docling:
  - pypdf is pure-Python, ~3MB, no native deps, ships in our env.
  - docling is ~400MB and pulls heavy native deps; not worth it for the
    typical pitch-deck ingestion use case (slide text, business plans, briefs).
  - If a user uploads a layout-heavy document where pypdf's text extraction
    is poor, we can swap in docling later behind the same interface.
"""

from __future__ import annotations

import json
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import structlog

logger = structlog.get_logger(__name__)


@dataclass
class ExtractedDocument:
    source: str
    mime_type: str
    pages: list[str] = field(default_factory=list)
    text: str = ""
    metadata: dict = field(default_factory=dict)
    error: Optional[str] = None

    @property
    def word_count(self) -> int:
        return len(self.text.split())


_TEXT_EXTS = {".txt", ".md", ".rst", ".log"}
_DATA_EXTS = {".csv", ".json"}


def extract_document(path: str | Path) -> ExtractedDocument:
    """Extract text from a single document. Sync — call via asyncio.to_thread."""
    p = Path(path)
    ext = p.suffix.lower()

    if not p.exists() or not p.is_file():
        return ExtractedDocument(source=str(p), mime_type="unknown", error="file not found")

    try:
        if ext == ".pdf":
            return _extract_pdf(p)
        if ext == ".docx":
            return _extract_docx(p)
        if ext == ".pptx":
            return _extract_pptx(p)
        if ext in _TEXT_EXTS:
            return _extract_plain_text(p)
        if ext in _DATA_EXTS:
            return _extract_data(p, ext)
        # Unknown extension — try plain text as a last resort
        return _extract_plain_text(p, mime="application/octet-stream")
    except Exception as e:
        logger.warning("doc_extract.failed", path=str(p), ext=ext, error=str(e))
        return ExtractedDocument(source=str(p), mime_type=ext.lstrip("."), error=str(e))


def _extract_pdf(p: Path) -> ExtractedDocument:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError as e:
        raise RuntimeError("pypdf not installed; add `pypdf>=4.0.0`") from e

    reader = PdfReader(str(p))
    pages: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            txt = page.extract_text() or ""
            pages.append(txt.strip())
        except Exception as e:
            logger.warning("doc_extract.pdf_page_failed", page=i, error=str(e))
            pages.append("")

    meta = {}
    try:
        info = reader.metadata or {}
        meta = {k.lstrip("/"): str(v) for k, v in info.items() if v is not None}
    except Exception:
        pass

    text = "\n\n".join(pp for pp in pages if pp).strip()
    return ExtractedDocument(
        source=str(p),
        mime_type="application/pdf",
        pages=pages,
        text=text,
        metadata={**meta, "n_pages": len(pages)},
    )


def _extract_docx(p: Path) -> ExtractedDocument:
    try:
        from docx import Document  # type: ignore
    except ImportError as e:
        raise RuntimeError("python-docx not installed; add `python-docx>=1.1.0`") from e

    doc = Document(str(p))
    paragraphs = [par.text.strip() for par in doc.paragraphs if par.text.strip()]
    # Tables → tab-separated rows
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                paragraphs.append("\t".join(cells))
    text = "\n".join(paragraphs).strip()
    return ExtractedDocument(
        source=str(p),
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        pages=[text],
        text=text,
        metadata={"n_paragraphs": len(paragraphs)},
    )


def _extract_pptx(p: Path) -> ExtractedDocument:
    """Extract text from PowerPoint presentations."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise RuntimeError("python-pptx not installed; add `python-pptx>=0.6.21`") from e

    prs = Presentation(str(p))
    slides_text: list[str] = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_content: list[str] = []
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_content.append(shape.text.strip())
        
        # Extract text from tables
        for shape in slide.shapes:
            if hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        slide_content.append("\t".join(cells))
        
        slide_text = "\n".join(slide_content).strip()
        if slide_text:
            slides_text.append(f"Slide {slide_idx + 1}:\n{slide_text}")
    
    text = "\n\n".join(slides_text).strip()
    return ExtractedDocument(
        source=str(p),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        pages=slides_text,
        text=text,
        metadata={"n_slides": len(prs.slides)},
    )


def extract_from_url(url: str) -> ExtractedDocument:
    """Extract text from a URL using web scraping."""
    try:
        from firecrawl import FirecrawlApp  # type: ignore
    except ImportError:
        # Fallback: try tavily if firecrawl not available
        try:
            import requests
            from bs4 import BeautifulSoup  # type: ignore
            
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, "html.parser")
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text(separator="\n", strip=True)
            # Clean up whitespace
            lines = [line.strip() for line in text.split("\n") if line.strip()]
            text = "\n".join(lines)
            
            return ExtractedDocument(
                source=url,
                mime_type="text/html",
                pages=[text],
                text=text,
                metadata={"scraper": "beautifulsoup", "url": url},
            )
        except ImportError:
            return ExtractedDocument(
                source=url,
                mime_type="text/html",
                error="Web scraping not available; install firecrawl or beautifulsoup"
            )
    
    # Use Firecrawl for better extraction
    app = FirecrawlApp()
    scrape_result = app.scrape_url(url, params={"formats": ["markdown"]})
    
    text = scrape_result.get("markdown", "") or scrape_result.get("content", "")
    return ExtractedDocument(
        source=url,
        mime_type="text/html",
        pages=[text],
        text=text.strip(),
        metadata={"scraper": "firecrawl", "url": url},
    )


def _extract_plain_text(p: Path, mime: str = "text/plain") -> ExtractedDocument:
    raw = p.read_text(encoding="utf-8", errors="replace")
    return ExtractedDocument(
        source=str(p),
        mime_type=mime,
        pages=[raw],
        text=raw.strip(),
        metadata={"n_chars": len(raw)},
    )


def _extract_data(p: Path, ext: str) -> ExtractedDocument:
    raw = p.read_text(encoding="utf-8", errors="replace")
    if ext == ".json":
        try:
            obj = json.loads(raw)
            text = json.dumps(obj, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            text = raw
        mime = "application/json"
    else:
        text = raw
        mime = "text/csv"
    return ExtractedDocument(
        source=str(p),
        mime_type=mime,
        pages=[text],
        text=text.strip(),
        metadata={"n_chars": len(text)},
    )


def chunk_text(text: str, *, chunk_chars: int = 1200, overlap: int = 150) -> list[str]:
    """Sliding-window chunker for embedding ingestion.

    Splits on paragraph boundaries when possible to keep chunks semantically clean.
    """
    if not text:
        return []
    if len(text) <= chunk_chars:
        return [text]
    chunks: list[str] = []
    cursor = 0
    while cursor < len(text):
        end = min(len(text), cursor + chunk_chars)
        # Try to end on a paragraph break for cleaner chunks
        if end < len(text):
            break_at = text.rfind("\n\n", cursor + chunk_chars // 2, end)
            if break_at > 0:
                end = break_at
        chunks.append(text[cursor:end].strip())
        if end >= len(text):
            break
        cursor = max(end - overlap, cursor + 1)
    return [c for c in chunks if c]
