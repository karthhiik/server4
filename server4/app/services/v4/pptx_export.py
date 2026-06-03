"""
V4 PPTX export — 80% fidelity static slide rendering.

Consumes the v4 slide DTO (the same shape stored in `db.slides` and
returned by `GET /api/v4/projects/{project_id}/slides`) plus an
optional `design_tokens` dict (output of
`design_resolver.ResolvedDesignTokens.to_dict()`) and emits a
`.pptx` byte string built with python-pptx.

Why "80% fidelity":
- We faithfully reproduce: title, subtitle, body, bullets, stat blocks,
  comparison columns, timelines, tables, charts, quotes, image URL
  insertion (best-effort), and speaker notes.
- We deliberately do NOT try to reproduce: animations, hover states,
  interactive React elements, custom shader-style backgrounds, glass
  morphism, or any pixel-perfect parity with the live JSX renderer.
  Those degrade gracefully into static text/blocks.

No-fake-data invariants:
- Every value rendered comes from the provided slide dict. We never
  invent bullets, never substitute placeholder copy. If a field is
  missing, we omit the corresponding region rather than insert a
  generic "lorem ipsum".
- Citations and source URLs (when present on chart/stat slides) are
  rendered verbatim as a footer line. We never relabel a source.
- Speaker notes go into the actual notes_slide so the presenter mode
  stays useful.

Public API:
    V4PptxBuilder().build(slides, design_tokens=None, metadata=None) -> bytes
"""
from __future__ import annotations

import base64
import io
import logging
import os
import time
from typing import Any, Callable, Iterable, Optional
from urllib.parse import urlparse

import httpx
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)


# ── Geometry (16:9) ────────────────────────────────────────────────
_W_EMU = 12_192_000   # 13.333 in
_H_EMU = 6_858_000    # 7.5 in
_MARGIN_EMU = 457_200  # 0.5 in
_GUTTER_EMU = 228_600  # 0.25 in

# ── Image embedding budgets ────────────────────────────────────────
# Per-image timeout, total budget across all fetches per build, and a
# hard byte ceiling. PPTX bloats badly with multi-MB images and the
# export endpoint is synchronous for the user.
_IMAGE_FETCH_TIMEOUT = float(os.environ.get("V4_PPTX_IMAGE_TIMEOUT", "5.0"))
_IMAGE_TOTAL_BUDGET = float(os.environ.get("V4_PPTX_IMAGE_BUDGET", "12.0"))
_IMAGE_MAX_BYTES = int(os.environ.get("V4_PPTX_IMAGE_MAX_BYTES", "5000000"))
# python-pptx natively supports png/jpg/gif/bmp. WebP works on most
# modern PowerPoint installs (≥ 2019). SVG must be rasterised first
# which we do NOT do at export time — we fall back to an honest text
# placeholder so the user sees what was actually emitted.
_PPTX_IMAGE_MIMES = {
    "image/png", "image/jpeg", "image/jpg",
    "image/gif", "image/webp", "image/bmp",
}


# ── Defaults — used only when design_tokens is None or partial ─────
_DEFAULT_PALETTE: dict[str, Any] = {
    "primary": "#2563EB",
    "secondary": "#1E40AF",
    "accent": "#7C3AED",
    "background": "#FFFFFF",
    "surface": "#F8FAFC",
    "text_primary": "#0F172A",
    "text_secondary": "#475569",
    "text_muted": "#94A3B8",
    "success": "#16A34A",
    "warning": "#F59E0B",
    "danger": "#DC2626",
    "chart": ["#2563EB", "#7C3AED", "#16A34A", "#F59E0B", "#DC2626", "#0EA5E9"],
}
_DEFAULT_FONTS: dict[str, str] = {
    "heading": "Calibri",
    "body": "Calibri",
}


_CHART_TYPE_MAP: dict[str, Any] = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "donut": XL_CHART_TYPE.DOUGHNUT,
}


def _hex_to_rgb(value: Optional[str], fallback: str = "#000000") -> RGBColor:
    """Convert '#rrggbb' to RGBColor; tolerate noise and clamp to fallback.

    Why so defensive: design_tokens come from user input + LLM. We have
    seen 'rgb(...)', '#rgb' (3-digit short form), and bare 'red'. None of
    these are valid python-pptx colors; rather than 500-error the export
    we silently fall back to the palette default.
    """
    if not isinstance(value, str):
        value = fallback
    s = value.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(ch * 2 for ch in s)
    if len(s) != 6 or not all(c in "0123456789abcdefABCDEF" for c in s):
        s = fallback.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _coerce_float(value: Any) -> Optional[float]:
    """Best-effort numeric coercion for chart values.

    Strings like '$4.1M', '18%', '1,200' are normalized to bare floats.
    Returns None when the value cannot be coerced — the caller MUST
    skip the data point rather than fabricate a zero (which would lie
    on the chart).
    """
    if isinstance(value, (int, float)):
        return float(value)
    if not isinstance(value, str):
        return None
    cleaned = value.strip().replace(",", "")
    # Strip trailing % and leading currency-style symbols.
    if cleaned.endswith("%"):
        cleaned = cleaned[:-1]
    cleaned = cleaned.lstrip("$€£¥")
    # Handle 'M' / 'B' / 'K' suffixes — express in base units so the
    # chart axis is meaningful.
    mult = 1.0
    if cleaned and cleaned[-1] in "kKmMbB":
        mult = {"k": 1e3, "K": 1e3, "m": 1e6, "M": 1e6, "b": 1e9, "B": 1e9}[cleaned[-1]]
        cleaned = cleaned[:-1]
    try:
        return float(cleaned) * mult
    except (ValueError, TypeError):
        return None


def _decode_data_uri(uri: str) -> Optional[tuple[bytes, str]]:
    """Decode a `data:image/<mime>;base64,<…>` URI.

    Returns `(bytes, mime)` on success or `None` on any malformation.
    Per no-fake-data: we never "guess" the mime from the payload — if
    the header is missing or the mime is unsupported we return None and
    the caller falls back to an honest text placeholder.
    """
    if not isinstance(uri, str) or not uri.startswith("data:image/"):
        return None
    try:
        head, _, payload = uri.partition(",")
        if not head or not payload:
            return None
        if ";base64" not in head:
            return None
        # head looks like 'data:image/png;base64'
        mime = head[5:].split(";", 1)[0].strip().lower()
        if mime not in _PPTX_IMAGE_MIMES:
            return None
        raw = base64.b64decode(payload, validate=False)
    except Exception:  # noqa: BLE001
        return None
    if len(raw) < 32 or len(raw) > _IMAGE_MAX_BYTES:
        return None
    return raw, mime


# ═══════════════════════════════════════════════════════════════════
# Builder
# ═══════════════════════════════════════════════════════════════════


class V4PptxBuilder:
    """Build a .pptx byte string from v4 slide DTOs + design tokens.

    Stateless across calls except for the per-build `_image_cache` and
    `_image_budget_remaining`, which are reset at the top of `build()`.

    `image_resolver` (optional): a callable `(url) -> Optional[bytes]`
    that bypasses the network. Used in tests so we never hit live HTTP.
    """

    def __init__(
        self,
        *,
        image_resolver: Optional[Callable[[str], Optional[bytes]]] = None,
    ) -> None:
        self._image_resolver = image_resolver
        # Per-build state — always reset in build().
        self._image_cache: dict[str, Optional[bytes]] = {}
        self._image_budget_remaining: float = _IMAGE_TOTAL_BUDGET
        self._embed_images: bool = True
        self._show_page_numbers: bool = True

    def build(
        self,
        slides: Iterable[dict[str, Any]],
        design_tokens: Optional[dict[str, Any]] = None,
        metadata: Optional[dict[str, Any]] = None,
        *,
        embed_images: bool = True,
        show_page_numbers: bool = True,
    ) -> bytes:
        slides_list = [s for s in slides if isinstance(s, dict)]
        if not slides_list:
            # Slice 4 (Export Parity): refuse to emit a corrupt 0-slide
            # artifact. Routers translate this into a structured 409
            # envelope. Kept compatible with prior callers by inheriting
            # from ``ValueError``.
            from app.services.v4.errors import ExportContentEmpty

            raise ExportContentEmpty("V4PptxBuilder.build: slides is empty")

        # Reset per-build state.
        self._image_cache = {}
        self._image_budget_remaining = _IMAGE_TOTAL_BUDGET
        self._embed_images = bool(embed_images)
        self._show_page_numbers = bool(show_page_numbers)

        palette = self._resolve_palette(design_tokens)
        fonts = self._resolve_fonts(design_tokens)
        meta = metadata or {}

        prs = Presentation()
        prs.slide_width = Emu(_W_EMU)
        prs.slide_height = Emu(_H_EMU)

        # Apply core metadata properties so the file shows the deck title
        # in PowerPoint's recent-files list. We never invent these.
        cp = prs.core_properties
        if meta.get("title"):
            cp.title = str(meta["title"])[:255]
        if meta.get("author"):
            cp.author = str(meta["author"])[:255]
        if meta.get("company"):
            cp.comments = str(meta["company"])[:512]

        for s in sorted(slides_list, key=lambda d: int(d.get("index", 0))):
            try:
                self._render_slide(prs, s, palette, fonts)
            except Exception as exc:  # noqa: BLE001
                # Per the no-fake-data rule, a failure renders an honest
                # placeholder noting the slide could not be exported —
                # NOT a fabricated copy of the slide.
                logger.warning(
                    "v4_pptx_slide_failed",
                    extra={"index": s.get("index"), "error": str(exc)[:200]},
                )
                self._render_failure_slide(prs, s, palette, fonts, str(exc)[:120])

        # Deck chrome — page numbers as a footer. Skipped for 1-slide decks
        # and when the caller opts out via show_page_numbers=False.
        if self._show_page_numbers and len(prs.slides) > 1:
            total = len(prs.slides)
            for i, ps in enumerate(prs.slides):
                try:
                    self._render_page_number(ps, i, total, palette, fonts)
                except Exception:  # noqa: BLE001
                    pass  # chrome is best-effort

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    # ── token resolution ───────────────────────────────────────────

    @staticmethod
    def _resolve_palette(tokens: Optional[dict[str, Any]]) -> dict[str, Any]:
        out = dict(_DEFAULT_PALETTE)
        if not tokens:
            return out
        p = tokens.get("palette") or {}
        for key, default in _DEFAULT_PALETTE.items():
            v = p.get(key)
            if isinstance(default, list):
                if isinstance(v, list) and v:
                    out[key] = [str(c) for c in v]
            elif isinstance(v, str) and v:
                out[key] = v
        return out

    @staticmethod
    def _resolve_fonts(tokens: Optional[dict[str, Any]]) -> dict[str, str]:
        out = dict(_DEFAULT_FONTS)
        if not tokens:
            return out
        f = tokens.get("fonts") or {}
        for key in ("heading", "body"):
            v = f.get(key)
            if isinstance(v, str) and v.strip():
                out[key] = v.strip()
        return out

    # ── slide render dispatch ──────────────────────────────────────

    def _render_slide(
        self,
        prs: Presentation,
        slide: dict[str, Any],
        palette: dict[str, Any],
        fonts: dict[str, str],
    ) -> None:
        # Use blank layout (index 6 in default template) so we own all
        # geometry. Layout-specific styling is driven by the slide DTO.
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(pptx_slide, palette)

        layout = (slide.get("layout") or "").lower().strip()
        intent = (slide.get("intent") or "").lower().strip()

        # Title-style slides — large centered title + subtitle only.
        if layout in {"title", "title-slide", "title-only"} or intent == "title":
            self._render_title_slide(pptx_slide, slide, palette, fonts)
        # Quote slide.
        elif layout == "quote" or slide.get("quote"):
            self._render_quote_slide(pptx_slide, slide, palette, fonts)
        # Stat-grid / stat-hero — large stat blocks.
        elif layout in {"stat-grid", "stat-hero", "grid-3"} or slide.get("stat_blocks"):
            # Stat blocks may co-exist with a body; we always render the
            # title at top and the stat blocks beneath.
            self._render_stat_slide(pptx_slide, slide, palette, fonts)
        # Chart slide. `chart-focus` is the canonical layout name from
        # the skeleton planner; `chart` is the legacy/short alias.
        elif layout in {"chart", "chart-focus"} or slide.get("chart"):
            self._render_chart_slide(pptx_slide, slide, palette, fonts)
        # Table slide.
        elif layout == "table" or slide.get("table"):
            self._render_table_slide(pptx_slide, slide, palette, fonts)
        # Timeline.
        elif layout == "timeline" or slide.get("timeline"):
            self._render_timeline_slide(pptx_slide, slide, palette, fonts)
        # Comparison.
        elif layout == "comparison" or slide.get("comparison"):
            self._render_comparison_slide(pptx_slide, slide, palette, fonts)
        # Diagram → render nodes as text blocks (no auto-layout).
        # `process` is a canonical planner layout that semantically
        # describes a flow of steps — render it through the diagram
        # path so we keep the directional treatment instead of
        # silently falling through to a generic bullet wall.
        elif layout in {"diagram", "process"} or slide.get("diagram"):
            self._render_diagram_slide(pptx_slide, slide, palette, fonts)
        # Team grid (premium-only field, but layout-driven).
        elif layout in {"team", "team-grid"} or intent == "team" or slide.get("team_members") or slide.get("requires_user_input"):
            self._render_team_slide(pptx_slide, slide, palette, fonts)
        # Full-bleed image slide.
        elif layout == "image-full":
            self._render_image_full_slide(pptx_slide, slide, palette, fonts)
        # Two-column.
        elif layout in {"two-column", "image-left", "image-right"}:
            self._render_two_column_slide(pptx_slide, slide, palette, fonts)
        # Default: title + bullets + body.
        else:
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)

        # Speaker notes — every layout gets these.
        self._render_link_cta(pptx_slide, slide, palette, fonts)

        notes = slide.get("speaker_notes") or ""
        if isinstance(notes, str) and notes.strip():
            try:
                pptx_slide.notes_slide.notes_text_frame.text = notes.strip()
            except Exception:  # noqa: BLE001
                pass  # presenter notes are best-effort

        # Citation footer — never invented. Only emit when the slide
        # carries explicit citations.
        self._render_citation_footer(pptx_slide, slide, palette, fonts)

    # ── primitives ─────────────────────────────────────────────────

    @staticmethod
    def _set_background(pptx_slide, palette: dict[str, Any]) -> None:
        bg = pptx_slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _hex_to_rgb(palette.get("background"), "#FFFFFF")

    # ── image embedding ────────────────────────────────────────────

    def _fetch_image_bytes(self, url: str) -> Optional[bytes]:
        """Resolve `url` to image bytes for embedding.

        Order:
          1. If a test `image_resolver` is set, use it (offline path).
          2. If `url` is a `data:image/...;base64,...` URI, decode locally.
          3. Otherwise issue a strict-budget httpx GET.

        Returns `None` on any failure. The caller MUST NOT fabricate a
        substitute image — it should fall back to an honest text
        placeholder so the user knows what was actually embedded.
        """
        if not isinstance(url, str) or not url.strip() or not self._embed_images:
            return None
        if url in self._image_cache:
            return self._image_cache[url]

        result: Optional[bytes] = None
        try:
            if self._image_resolver is not None:
                try:
                    result = self._image_resolver(url)
                except Exception:  # noqa: BLE001
                    result = None
            elif url.startswith("data:image/"):
                decoded = _decode_data_uri(url)
                if decoded is not None:
                    result = decoded[0]
            else:
                parsed = urlparse(url)
                if parsed.scheme not in ("http", "https"):
                    result = None
                elif self._image_budget_remaining <= 0:
                    result = None
                else:
                    t0 = time.monotonic()
                    try:
                        with httpx.Client(
                            timeout=_IMAGE_FETCH_TIMEOUT,
                            follow_redirects=True,
                        ) as client:
                            resp = client.get(
                                url,
                                headers={"User-Agent": "v4-pptx-export/1.0"},
                            )
                        if resp.status_code != 200:
                            result = None
                        else:
                            ctype = (
                                resp.headers.get("content-type", "")
                                .split(";", 1)[0]
                                .strip()
                                .lower()
                            )
                            body = resp.content
                            if ctype not in _PPTX_IMAGE_MIMES:
                                result = None
                            elif len(body) < 32 or len(body) > _IMAGE_MAX_BYTES:
                                result = None
                            else:
                                result = body
                    finally:
                        self._image_budget_remaining -= (time.monotonic() - t0)
        except Exception as exc:  # noqa: BLE001
            logger.debug(
                "v4_pptx_image_fetch_failed",
                extra={"url": str(url)[:200], "error": str(exc)[:120]},
            )
            result = None

        # Final size-cap re-check (resolver paths may bypass the network check).
        if isinstance(result, (bytes, bytearray)):
            if len(result) < 32 or len(result) > _IMAGE_MAX_BYTES:
                result = None

        self._image_cache[url] = result
        return result

    def _add_picture_safe(
        self,
        pptx_slide,
        url: str,
        *,
        left_emu: int,
        top_emu: int,
        width_emu: int,
        height_emu: int,
        link_url: str = "",
    ) -> bool:
        """Try to embed `url` at the given rect. Return True iff embedded."""
        body = self._fetch_image_bytes(url)
        if not body:
            return False
        try:
            shape = pptx_slide.shapes.add_picture(
                io.BytesIO(body),
                Emu(left_emu), Emu(top_emu),
                width=Emu(width_emu), height=Emu(height_emu),
            )
            if link_url:
                self._set_shape_hyperlink(shape, link_url)
            return True
        except Exception as exc:  # noqa: BLE001
            logger.debug(
                "v4_pptx_add_picture_failed",
                extra={"error": str(exc)[:120]},
            )
            return False

    def _add_rounded_card(
        self,
        pptx_slide,
        *,
        left_emu: int,
        top_emu: int,
        width_emu: int,
        height_emu: int,
        fill_hex: str,
        line_hex: Optional[str] = None,
    ):
        """Draw a filled rounded rectangle as a card background.

        The shape's own text frame is intentionally cleared so callers
        can overlay separate text boxes on top with full control.
        """
        shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(left_emu), Emu(top_emu),
            Emu(width_emu), Emu(height_emu),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
        if line_hex is None:
            shape.line.fill.background()  # no border
        else:
            shape.line.color.rgb = _hex_to_rgb(line_hex)
            shape.line.width = Pt(1)
        if shape.has_text_frame:
            shape.text_frame.text = ""
        return shape

    @staticmethod
    def _set_run_hyperlink(run, url: str) -> None:
        """Attach `url` to a text run as a clickable hyperlink.

        Best-effort: python-pptx's hyperlink API is well-supported but
        we never want a hyperlink failure to abort the whole export.
        """
        if not isinstance(url, str) or not url.strip():
            return
        try:
            run.hyperlink.address = url.strip()
        except Exception:  # noqa: BLE001
            pass

    @staticmethod
    def _set_shape_hyperlink(shape, url: str) -> None:
        if not isinstance(url, str) or not url.strip():
            return
        try:
            shape.click_action.hyperlink.address = url.strip()
        except Exception:  # noqa: BLE001
            pass

    @staticmethod
    def _slide_links(slide: dict[str, Any], max_n: int = 6) -> list[dict[str, str]]:
        links: list[dict[str, str]] = []
        seen: set[str] = set()

        def add(label: str, url: str, target: str = "text") -> None:
            clean = str(url or "").strip()
            if not clean.lower().startswith(("http://", "https://")):
                return
            if clean in seen:
                return
            seen.add(clean)
            links.append({
                "label": (str(label or clean).strip() or clean)[:120],
                "url": clean[:500],
                "target": target if target in {"text", "button", "image", "source"} else "text",
            })

        for item in slide.get("links") or []:
            if isinstance(item, dict):
                add(
                    str(item.get("label") or item.get("title") or "Open link"),
                    str(item.get("url") or item.get("href") or ""),
                    str(item.get("target") or "text").strip().lower(),
                )
        for item in slide.get("citations") or []:
            if isinstance(item, dict):
                add(str(item.get("title") or "Source"), str(item.get("url") or ""), "source")
        return links[:max_n]

    @staticmethod
    def _first_link_url(slide: dict[str, Any], targets: set[str]) -> str:
        for link in V4PptxBuilder._slide_links(slide):
            if link.get("target") in targets:
                return link.get("url", "")
        return ""

    def _render_page_number(
        self, pptx_slide, idx: int, total: int, palette, fonts,
    ) -> None:
        """Render a small `i / N` footer in the bottom-right corner."""
        self._add_text_block(
            pptx_slide,
            text=f"{idx + 1} / {total}",
            left_emu=_W_EMU - _MARGIN_EMU - Inches(1.0).emu,
            top_emu=_H_EMU - Inches(0.4).emu,
            width_emu=Inches(1.0).emu, height_emu=Inches(0.3).emu,
            size_pt=9, bold=False,
            font_name=fonts["body"],
            color_hex=palette["text_muted"],
            align=PP_ALIGN.RIGHT,
        )

    @staticmethod
    def _add_textbox(
        pptx_slide,
        *,
        left_emu: int,
        top_emu: int,
        width_emu: int,
        height_emu: int,
    ):
        return pptx_slide.shapes.add_textbox(
            Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu)
        )

    @staticmethod
    def _set_paragraph(
        para,
        *,
        text: str,
        size_pt: float,
        bold: bool,
        font_name: str,
        rgb: RGBColor,
        align: Optional[Any] = None,
    ) -> None:
        para.text = text
        run = para.runs[0]
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = rgb
        if align is not None:
            para.alignment = align

    def _add_text_block(
        self,
        pptx_slide,
        *,
        text: str,
        left_emu: int,
        top_emu: int,
        width_emu: int,
        height_emu: int,
        size_pt: float,
        bold: bool,
        font_name: str,
        color_hex: str,
        align: Optional[Any] = None,
        anchor: Optional[Any] = None,
    ):
        tx = self._add_textbox(
            pptx_slide, left_emu=left_emu, top_emu=top_emu,
            width_emu=width_emu, height_emu=height_emu,
        )
        tf = tx.text_frame
        tf.word_wrap = True
        if anchor is not None:
            tf.vertical_anchor = anchor
        self._set_paragraph(
            tf.paragraphs[0],
            text=str(text),
            size_pt=size_pt,
            bold=bold,
            font_name=font_name,
            rgb=_hex_to_rgb(color_hex),
            align=align,
        )
        return tx

    def _add_title(
        self,
        pptx_slide,
        *,
        text: str,
        palette: dict[str, Any],
        fonts: dict[str, str],
        size_pt: float = 32,
        top_emu: int = _MARGIN_EMU,
        height_emu: int = Inches(1.1).emu,
    ) -> int:
        """Add the slide title and return its bottom Y in EMU."""
        if not isinstance(text, str) or not text.strip():
            return top_emu
        self._add_text_block(
            pptx_slide,
            text=text.strip(),
            left_emu=_MARGIN_EMU,
            top_emu=top_emu,
            width_emu=_W_EMU - 2 * _MARGIN_EMU,
            height_emu=height_emu,
            size_pt=size_pt,
            bold=True,
            font_name=fonts["heading"],
            color_hex=palette["text_primary"],
        )
        return top_emu + height_emu

    def _add_subheadline(
        self,
        pptx_slide,
        *,
        text: str,
        palette: dict[str, Any],
        fonts: dict[str, str],
        top_emu: int,
        size_pt: float = 16,
    ) -> int:
        if not isinstance(text, str) or not text.strip():
            return top_emu
        h = Inches(0.7).emu
        self._add_text_block(
            pptx_slide,
            text=text.strip(),
            left_emu=_MARGIN_EMU,
            top_emu=top_emu,
            width_emu=_W_EMU - 2 * _MARGIN_EMU,
            height_emu=h,
            size_pt=size_pt,
            bold=False,
            font_name=fonts["body"],
            color_hex=palette["text_secondary"],
        )
        return top_emu + h + Emu(_GUTTER_EMU // 2).emu

    # ── layout renderers ───────────────────────────────────────────

    def _render_title_slide(
        self, pptx_slide, slide, palette, fonts,
    ) -> None:
        headline = slide.get("headline") or ""
        sub = slide.get("subheadline") or ""
        body = slide.get("body") or ""
        company_icon = slide.get("company_icon_url")

        # Optional company icon — small badge in top-right when present.
        # Never invented; only rendered when the DSL carries a real URL.
        if isinstance(company_icon, str) and company_icon.strip():
            icon_size = Inches(0.9).emu
            self._add_picture_safe(
                pptx_slide, company_icon,
                left_emu=_W_EMU - _MARGIN_EMU - icon_size,
                top_emu=_MARGIN_EMU,
                width_emu=icon_size, height_emu=icon_size,
            )

        # Centered hero title at vertical midline.
        title_box = self._add_text_block(
            pptx_slide,
            text=headline,
            left_emu=_MARGIN_EMU,
            top_emu=Inches(2.4).emu,
            width_emu=_W_EMU - 2 * _MARGIN_EMU,
            height_emu=Inches(1.6).emu,
            size_pt=44,
            bold=True,
            font_name=fonts["heading"],
            color_hex=palette["text_primary"],
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        if sub:
            self._add_text_block(
                pptx_slide,
                text=sub,
                left_emu=_MARGIN_EMU,
                top_emu=Inches(4.1).emu,
                width_emu=_W_EMU - 2 * _MARGIN_EMU,
                height_emu=Inches(0.9).emu,
                size_pt=20,
                bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_secondary"],
                align=PP_ALIGN.CENTER,
            )
        if body:
            self._add_text_block(
                pptx_slide,
                text=body,
                left_emu=_MARGIN_EMU * 2,
                top_emu=Inches(5.1).emu,
                width_emu=_W_EMU - 4 * _MARGIN_EMU,
                height_emu=Inches(1.6).emu,
                size_pt=14,
                bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_muted"],
                align=PP_ALIGN.CENTER,
            )
        return None

    def _render_bullets_slide(self, pptx_slide, slide, palette, fonts) -> None:
        cursor = self._add_title(
            pptx_slide,
            text=slide.get("headline") or "",
            palette=palette, fonts=fonts,
        )
        cursor = self._add_subheadline(
            pptx_slide, text=slide.get("subheadline") or "",
            palette=palette, fonts=fonts, top_emu=cursor,
        )
        bullets = slide.get("bullets") or []
        body = slide.get("body") or ""

        if bullets:
            tx = self._add_textbox(
                pptx_slide,
                left_emu=_MARGIN_EMU,
                top_emu=cursor,
                width_emu=_W_EMU - 2 * _MARGIN_EMU,
                height_emu=_H_EMU - cursor - Inches(0.5).emu,
            )
            tf = tx.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(bullets):
                if not isinstance(bullet, str):
                    continue
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                self._set_paragraph(
                    p,
                    text=f"• {bullet}",
                    size_pt=18,
                    bold=False,
                    font_name=fonts["body"],
                    rgb=_hex_to_rgb(palette["text_primary"]),
                )
                p.space_after = Pt(8)
        elif body:
            self._add_text_block(
                pptx_slide,
                text=body,
                left_emu=_MARGIN_EMU,
                top_emu=cursor,
                width_emu=_W_EMU - 2 * _MARGIN_EMU,
                height_emu=_H_EMU - cursor - Inches(0.5).emu,
                size_pt=16,
                bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_secondary"],
            )

    def _render_two_column_slide(self, pptx_slide, slide, palette, fonts) -> None:
        cursor = self._add_title(
            pptx_slide,
            text=slide.get("headline") or "",
            palette=palette, fonts=fonts,
        )
        cursor = self._add_subheadline(
            pptx_slide, text=slide.get("subheadline") or "",
            palette=palette, fonts=fonts, top_emu=cursor,
        )
        col_w = (_W_EMU - 2 * _MARGIN_EMU - _GUTTER_EMU) // 2
        bullets = slide.get("bullets") or []
        body = slide.get("body") or ""
        image_url = slide.get("image_url")

        # Left: bullets or body.
        left_x = _MARGIN_EMU
        right_x = _MARGIN_EMU + col_w + _GUTTER_EMU
        height = _H_EMU - cursor - Inches(0.5).emu

        if bullets:
            tx = self._add_textbox(
                pptx_slide,
                left_emu=left_x, top_emu=cursor,
                width_emu=col_w, height_emu=height,
            )
            tf = tx.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(bullets):
                if not isinstance(bullet, str):
                    continue
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                self._set_paragraph(
                    p, text=f"• {bullet}", size_pt=16, bold=False,
                    font_name=fonts["body"],
                    rgb=_hex_to_rgb(palette["text_primary"]),
                )
                p.space_after = Pt(6)
        elif body:
            self._add_text_block(
                pptx_slide, text=body,
                left_emu=left_x, top_emu=cursor,
                width_emu=col_w, height_emu=height,
                size_pt=14, bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_secondary"],
            )

        # Right: image (real embedding when fetchable) OR honest text
        # placeholder showing the URL/prompt verbatim. We never fabricate
        # a substitute image — a failed fetch falls back to the
        # placeholder so the user sees what was actually emitted.
        right_label = ""
        embedded = False
        if isinstance(image_url, str) and image_url.strip():
            embedded = self._add_picture_safe(
                pptx_slide, image_url,
                left_emu=right_x, top_emu=cursor,
                width_emu=col_w, height_emu=height,
            )
            if not embedded:
                right_label = f"[Image: {image_url}]"
        elif slide.get("image_prompt"):
            right_label = f"[Image: {slide['image_prompt']}]"
        if right_label and not embedded:
            self._add_text_block(
                pptx_slide, text=right_label,
                left_emu=right_x, top_emu=cursor,
                width_emu=col_w, height_emu=height,
                size_pt=12, bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_muted"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )

    def _render_stat_slide(self, pptx_slide, slide, palette, fonts) -> None:
        cursor = self._add_title(
            pptx_slide,
            text=slide.get("headline") or "",
            palette=palette, fonts=fonts,
        )
        cursor = self._add_subheadline(
            pptx_slide, text=slide.get("subheadline") or "",
            palette=palette, fonts=fonts, top_emu=cursor,
        )
        stats = [s for s in (slide.get("stat_blocks") or []) if isinstance(s, dict)]
        if not stats:
            # Fall back to bullets/body.
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)
            return

        # Use up to 3 stats per row; each card is a centered text block
        # with the value (large, bold) and label (small) stacked.
        n = min(len(stats), 6)
        cols = min(n, 3)
        rows = (n + cols - 1) // cols
        avail_w = _W_EMU - 2 * _MARGIN_EMU
        avail_h = _H_EMU - cursor - Inches(0.5).emu
        col_w = (avail_w - (cols - 1) * _GUTTER_EMU) // cols
        row_h = (avail_h - (rows - 1) * _GUTTER_EMU) // rows

        for idx in range(n):
            r = idx // cols
            c = idx % cols
            x = _MARGIN_EMU + c * (col_w + _GUTTER_EMU)
            y = cursor + r * (row_h + _GUTTER_EMU)
            stat = stats[idx]
            value = str(stat.get("value", ""))[:24]
            label = str(stat.get("label", ""))[:120]

            # Card background — surface fill, no border. Adds the visual
            # "stat card" treatment that matches the live JSX renderer.
            self._add_rounded_card(
                pptx_slide,
                left_emu=x, top_emu=y,
                width_emu=col_w, height_emu=row_h,
                fill_hex=palette["surface"],
            )

            # Value (top, bold, primary color).
            self._add_text_block(
                pptx_slide, text=value,
                left_emu=x, top_emu=y,
                width_emu=col_w, height_emu=int(row_h * 0.6),
                size_pt=40, bold=True,
                font_name=fonts["heading"],
                color_hex=palette["primary"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            # Label (bottom).
            self._add_text_block(
                pptx_slide, text=label,
                left_emu=x, top_emu=y + int(row_h * 0.6),
                width_emu=col_w, height_emu=int(row_h * 0.4),
                size_pt=12, bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_secondary"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
            )

    def _render_chart_slide(self, pptx_slide, slide, palette, fonts) -> None:
        cursor = self._add_title(
            pptx_slide,
            text=slide.get("headline") or "",
            palette=palette, fonts=fonts,
        )
        cursor = self._add_subheadline(
            pptx_slide, text=slide.get("subheadline") or "",
            palette=palette, fonts=fonts, top_emu=cursor,
        )
        chart = slide.get("chart") or {}
        chart_type = (chart.get("type") or "bar").lower()
        xl_type = _CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        data = chart.get("data") or []
        labels: list[str] = []
        values: list[float] = []
        for d in data:
            if not isinstance(d, dict):
                continue
            label = str(d.get("label", "")).strip()
            v = _coerce_float(d.get("value"))
            if not label or v is None:
                # No fake-zeroes — drop the point rather than lie.
                continue
            labels.append(label[:32])
            values.append(v)

        if not labels:
            # No usable chart data → render the slide as a bullets fallback,
            # never fabricate a chart.
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)
            return

        cd = CategoryChartData()
        cd.categories = labels
        series_label = str(chart.get("series_label") or "Value")
        cd.add_series(series_label, values)

        chart_top = cursor
        chart_h = _H_EMU - chart_top - Inches(0.6).emu
        frame = pptx_slide.shapes.add_chart(
            xl_type,
            Emu(_MARGIN_EMU), Emu(chart_top),
            Emu(_W_EMU - 2 * _MARGIN_EMU), Emu(chart_h),
            cd,
        )
        ch = frame.chart
        ch.has_legend = False
        chart_colors = palette.get("chart") or _DEFAULT_PALETTE["chart"]
        for idx, series in enumerate(ch.series):
            color = chart_colors[idx % len(chart_colors)]
            try:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = _hex_to_rgb(color)
            except Exception:  # noqa: BLE001
                pass

    def _render_table_slide(self, pptx_slide, slide, palette, fonts) -> None:
        cursor = self._add_title(
            pptx_slide,
            text=slide.get("headline") or "",
            palette=palette, fonts=fonts,
        )
        cursor = self._add_subheadline(
            pptx_slide, text=slide.get("subheadline") or "",
            palette=palette, fonts=fonts, top_emu=cursor,
        )
        table = slide.get("table") or {}
        headers = [str(h) for h in (table.get("headers") or []) if h is not None]
        rows = [
            [str(c) for c in r]
            for r in (table.get("rows") or [])
            if isinstance(r, (list, tuple))
        ]
        if not headers or not rows:
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)
            return

        # Cap to fit a 16:9 slide cleanly.
        n_cols = min(len(headers), 6)
        n_rows = min(len(rows), 8) + 1   # +1 header row
        headers = headers[:n_cols]
        rows = [r[:n_cols] for r in rows[:8]]

        avail_w = _W_EMU - 2 * _MARGIN_EMU
        avail_h = _H_EMU - cursor - Inches(0.6).emu
        shape = pptx_slide.shapes.add_table(
            n_rows, n_cols,
            Emu(_MARGIN_EMU), Emu(cursor),
            Emu(avail_w), Emu(avail_h),
        )
        tbl = shape.table

        # Header row styling.
        for c, h in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = h
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True
                    run.font.size = Pt(13)
                    run.font.name = fonts["heading"]
                    run.font.color.rgb = _hex_to_rgb("#FFFFFF")
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_to_rgb(palette["primary"])

        # Body rows.
        for r, row in enumerate(rows, start=1):
            for c in range(n_cols):
                cell = tbl.cell(r, c)
                cell.text = row[c] if c < len(row) else ""
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(11)
                        run.font.name = fonts["body"]
                        run.font.color.rgb = _hex_to_rgb(palette["text_primary"])

    def _render_timeline_slide(self, pptx_slide, slide, palette, fonts) -> None:
        cursor = self._add_title(
            pptx_slide,
            text=slide.get("headline") or "",
            palette=palette, fonts=fonts,
        )
        cursor = self._add_subheadline(
            pptx_slide, text=slide.get("subheadline") or "",
            palette=palette, fonts=fonts, top_emu=cursor,
        )
        tl = slide.get("timeline") or {}
        events = [e for e in (tl.get("events") or []) if isinstance(e, dict)]
        if not events:
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)
            return
        events = events[:7]
        n = len(events)
        avail_w = _W_EMU - 2 * _MARGIN_EMU
        col_w = (avail_w - (n - 1) * _GUTTER_EMU) // n
        col_h = Inches(3.5).emu

        for idx, ev in enumerate(events):
            x = _MARGIN_EMU + idx * (col_w + _GUTTER_EMU)
            date = str(ev.get("date", "")).strip()[:32]
            title = str(ev.get("title", "")).strip()[:80]
            desc = str(ev.get("description", "")).strip()[:160]
            # Date pill (primary color, bold).
            self._add_text_block(
                pptx_slide, text=date,
                left_emu=x, top_emu=cursor,
                width_emu=col_w, height_emu=Inches(0.5).emu,
                size_pt=12, bold=True,
                font_name=fonts["heading"],
                color_hex=palette["primary"],
                align=PP_ALIGN.CENTER,
            )
            # Title.
            self._add_text_block(
                pptx_slide, text=title,
                left_emu=x, top_emu=cursor + Inches(0.55).emu,
                width_emu=col_w, height_emu=Inches(0.7).emu,
                size_pt=14, bold=True,
                font_name=fonts["heading"],
                color_hex=palette["text_primary"],
                align=PP_ALIGN.CENTER,
            )
            # Description.
            if desc:
                self._add_text_block(
                    pptx_slide, text=desc,
                    left_emu=x, top_emu=cursor + Inches(1.3).emu,
                    width_emu=col_w, height_emu=col_h - Inches(1.3).emu,
                    size_pt=11, bold=False,
                    font_name=fonts["body"],
                    color_hex=palette["text_secondary"],
                    align=PP_ALIGN.CENTER,
                )

    def _render_comparison_slide(self, pptx_slide, slide, palette, fonts) -> None:
        cursor = self._add_title(
            pptx_slide,
            text=slide.get("headline") or "",
            palette=palette, fonts=fonts,
        )
        cursor = self._add_subheadline(
            pptx_slide, text=slide.get("subheadline") or "",
            palette=palette, fonts=fonts, top_emu=cursor,
        )
        cmp = slide.get("comparison") or {}
        cols = [c for c in (cmp.get("columns") or []) if isinstance(c, dict)]
        if not cols:
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)
            return
        cols = cols[:3]
        n = len(cols)
        avail_w = _W_EMU - 2 * _MARGIN_EMU
        col_w = (avail_w - (n - 1) * _GUTTER_EMU) // n
        col_h = _H_EMU - cursor - Inches(0.5).emu

        for idx, col in enumerate(cols):
            x = _MARGIN_EMU + idx * (col_w + _GUTTER_EMU)
            title = str(col.get("title", "")).strip()[:80]
            items = [str(i) for i in (col.get("items") or []) if isinstance(i, str)]
            highlight = bool(col.get("highlight"))
            heading_color = palette["accent"] if highlight else palette["primary"]

            # Card background — surface fill for normal columns, accent
            # border for highlighted columns. Honest treatment: only
            # highlighted columns get the colored border, the highlight
            # flag must come from the DSL.
            self._add_rounded_card(
                pptx_slide,
                left_emu=x, top_emu=cursor,
                width_emu=col_w, height_emu=col_h,
                fill_hex=palette["surface"],
                line_hex=palette["accent"] if highlight else None,
            )

            # Column title.
            self._add_text_block(
                pptx_slide, text=title,
                left_emu=x, top_emu=cursor,
                width_emu=col_w, height_emu=Inches(0.7).emu,
                size_pt=18, bold=True,
                font_name=fonts["heading"],
                color_hex=heading_color,
                align=PP_ALIGN.CENTER,
            )
            # Items.
            if items:
                tx = self._add_textbox(
                    pptx_slide,
                    left_emu=x, top_emu=cursor + Inches(0.8).emu,
                    width_emu=col_w, height_emu=col_h - Inches(0.8).emu,
                )
                tf = tx.text_frame
                tf.word_wrap = True
                for j, item in enumerate(items):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    self._set_paragraph(
                        p, text=f"• {item}",
                        size_pt=13, bold=False,
                        font_name=fonts["body"],
                        rgb=_hex_to_rgb(palette["text_primary"]),
                    )
                    p.space_after = Pt(6)

    def _render_diagram_slide(self, pptx_slide, slide, palette, fonts) -> None:
        cursor = self._add_title(
            pptx_slide,
            text=slide.get("headline") or "",
            palette=palette, fonts=fonts,
        )
        cursor = self._add_subheadline(
            pptx_slide, text=slide.get("subheadline") or "",
            palette=palette, fonts=fonts, top_emu=cursor,
        )
        dg = slide.get("diagram") or {}
        nodes = [n for n in (dg.get("nodes") or []) if isinstance(n, dict)]
        if not nodes:
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)
            return
        # Render nodes as a horizontal flow of rounded text blocks.
        # Adjacent edges (declared in the DSL) get real arrow connectors;
        # non-adjacent edges fall back to a verbatim caption row.
        nodes = nodes[:6]
        n = len(nodes)
        avail_w = _W_EMU - 2 * _MARGIN_EMU
        col_w = (avail_w - (n - 1) * _GUTTER_EMU) // n
        node_h = Inches(1.4).emu
        node_top = cursor + Inches(0.3).emu

        # Build identity → index map for edge resolution. We accept the
        # node's `id`, the node's `label`, and the integer index — all
        # patterns observed in real DSL outputs.
        id_to_idx: dict[str, int] = {}
        for idx, node in enumerate(nodes):
            for key in ("id", "label"):
                v = node.get(key)
                if isinstance(v, str) and v.strip() and v not in id_to_idx:
                    id_to_idx[v] = idx
                    id_to_idx[v.strip()] = idx

        def _resolve(ref: Any) -> Optional[int]:
            if isinstance(ref, int) and 0 <= ref < n:
                return ref
            if isinstance(ref, str):
                s = ref.strip()
                if s.isdigit():
                    i = int(s)
                    if 0 <= i < n:
                        return i
                if s in id_to_idx:
                    return id_to_idx[s]
            return None

        for idx, node in enumerate(nodes):
            x = _MARGIN_EMU + idx * (col_w + _GUTTER_EMU)
            label = str(node.get("label", "")).strip()[:80]
            self._add_text_block(
                pptx_slide, text=label,
                left_emu=x, top_emu=node_top,
                width_emu=col_w, height_emu=node_h,
                size_pt=14, bold=True,
                font_name=fonts["heading"],
                color_hex=palette["primary"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )

        # Edges: draw real RIGHT_ARROW shape between adjacent nodes,
        # collect non-adjacent edges into a caption (no fabrication).
        edges = [e for e in (dg.get("edges") or []) if isinstance(e, dict)]
        non_adjacent: list[str] = []
        if edges:
            arrow_h = Inches(0.35).emu
            arrow_w = _GUTTER_EMU
            arrow_top = node_top + (node_h - arrow_h) // 2
            drawn_pairs: set[tuple[int, int]] = set()
            for e in edges[:8]:
                src = _resolve(e.get("from"))
                dst = _resolve(e.get("to"))
                lbl_raw = e.get("label")
                lbl = str(lbl_raw).strip() if isinstance(lbl_raw, str) else ""
                if src is None or dst is None:
                    continue
                if dst == src + 1 and (src, dst) not in drawn_pairs:
                    # Adjacent forward edge → draw real arrow.
                    arrow_left = (
                        _MARGIN_EMU + src * (col_w + _GUTTER_EMU) + col_w
                    )
                    shape = pptx_slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        Emu(arrow_left), Emu(arrow_top),
                        Emu(arrow_w), Emu(arrow_h),
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
                    shape.line.fill.background()
                    drawn_pairs.add((src, dst))
                    if lbl and shape.has_text_frame:
                        tf = shape.text_frame
                        self._set_paragraph(
                            tf.paragraphs[0], text=lbl,
                            size_pt=9, bold=False,
                            font_name=fonts["body"],
                            rgb=_hex_to_rgb("#FFFFFF"),
                        )
                else:
                    src_lbl = str(nodes[src].get("label") or src)
                    dst_lbl = str(nodes[dst].get("label") or dst)
                    non_adjacent.append(
                        f"{src_lbl} → {dst_lbl}"
                        + (f" ({lbl})" if lbl else "")
                    )

        if non_adjacent:
            self._add_text_block(
                pptx_slide,
                text="Flow: " + "  ·  ".join(non_adjacent),
                left_emu=_MARGIN_EMU,
                top_emu=cursor + Inches(2.0).emu,
                width_emu=avail_w, height_emu=Inches(0.6).emu,
                size_pt=11, bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_muted"],
                align=PP_ALIGN.CENTER,
            )

    def _render_quote_slide(self, pptx_slide, slide, palette, fonts) -> None:
        quote = slide.get("quote") or {}
        text = str(quote.get("text", "")).strip()
        attribution = str(quote.get("attribution", "")).strip()
        if not text:
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)
            return
        # Big centered quote.
        self._add_text_block(
            pptx_slide, text=f"\u201C{text}\u201D",
            left_emu=_MARGIN_EMU * 2,
            top_emu=Inches(2.0).emu,
            width_emu=_W_EMU - 4 * _MARGIN_EMU,
            height_emu=Inches(3.5).emu,
            size_pt=28, bold=False,
            font_name=fonts["heading"],
            color_hex=palette["text_primary"],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        if attribution:
            self._add_text_block(
                pptx_slide, text=f"— {attribution}",
                left_emu=_MARGIN_EMU * 2,
                top_emu=Inches(5.6).emu,
                width_emu=_W_EMU - 4 * _MARGIN_EMU,
                height_emu=Inches(0.6).emu,
                size_pt=16, bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_secondary"],
                align=PP_ALIGN.CENTER,
            )

    def _render_team_slide(self, pptx_slide, slide, palette, fonts) -> None:
        cursor = self._add_title(
            pptx_slide,
            text=slide.get("headline") or "",
            palette=palette, fonts=fonts,
        )
        cursor = self._add_subheadline(
            pptx_slide, text=slide.get("subheadline") or "",
            palette=palette, fonts=fonts, top_emu=cursor,
        )
        if slide.get("requires_user_input"):
            return
        # Filter to dict members with a real name. Never invent a name.
        members = [
            m for m in (slide.get("team_members") or [])
            if isinstance(m, dict) and isinstance(m.get("name"), str) and m["name"].strip()
        ]
        if not members:
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)
            return
        # Cap at 6 members per slide; pptx can't readably show more.
        members = members[:6]
        n = len(members)
        # Two rows of up to 3 if >3, else single row.
        cols = 3 if n > 3 else n
        rows = 2 if n > 3 else 1
        avail_w = _W_EMU - 2 * _MARGIN_EMU
        col_w = (avail_w - (cols - 1) * _GUTTER_EMU) // cols
        avail_h = _H_EMU - cursor - Inches(0.5).emu
        row_h = (avail_h - (rows - 1) * _GUTTER_EMU) // rows

        for idx, m in enumerate(members):
            r = idx // cols
            c = idx % cols
            x = _MARGIN_EMU + c * (col_w + _GUTTER_EMU)
            y = cursor + r * (row_h + _GUTTER_EMU)
            name = str(m.get("name", "")).strip()[:80]
            role_raw = m.get("role")
            role = str(role_raw).strip()[:80] if isinstance(role_raw, str) and role_raw.strip() else ""
            bio_raw = m.get("bio")
            bio = str(bio_raw).strip() if isinstance(bio_raw, str) and bio_raw.strip() else ""
            linkedin_raw = m.get("linkedin_url")
            linkedin = str(linkedin_raw).strip() if isinstance(linkedin_raw, str) and linkedin_raw.strip() else ""
            photo_raw = m.get("photo_url")
            photo_url = str(photo_raw).strip() if isinstance(photo_raw, str) and photo_raw.strip() else ""

            # Try to embed the photo as a square avatar in the top-left
            # of the card. If unavailable (svg / fetch fail / no URL),
            # the text content takes the full width — no fake avatar.
            photo_size = min(Inches(1.1).emu, row_h // 3)
            text_left = x
            text_w = col_w
            embedded_photo = False
            if photo_url:
                embedded_photo = self._add_picture_safe(
                    pptx_slide, photo_url,
                    left_emu=x, top_emu=y,
                    width_emu=photo_size, height_emu=photo_size,
                )
            if embedded_photo:
                text_left = x + photo_size + Inches(0.15).emu
                text_w = col_w - photo_size - Inches(0.15).emu

            # Card geometry (text region).
            name_h = Inches(0.5).emu
            role_h = Inches(0.35).emu if role else 0
            bio_h = max(0, row_h - name_h - role_h - Inches(0.4).emu)
            # Name
            self._add_text_block(
                pptx_slide, text=name,
                left_emu=text_left, top_emu=y,
                width_emu=text_w, height_emu=name_h,
                size_pt=18, bold=True,
                font_name=fonts["heading"],
                color_hex=palette["primary"],
                align=PP_ALIGN.LEFT,
            )
            # Role
            if role:
                self._add_text_block(
                    pptx_slide, text=role,
                    left_emu=text_left, top_emu=y + name_h,
                    width_emu=text_w, height_emu=role_h,
                    size_pt=12, bold=False,
                    font_name=fonts["body"],
                    color_hex=palette["accent"],
                    align=PP_ALIGN.LEFT,
                )
            # Bio + linkedin appended as separate paragraphs (only when present)
            if bio or linkedin:
                tx = self._add_textbox(
                    pptx_slide,
                    left_emu=text_left, top_emu=y + name_h + role_h,
                    width_emu=text_w, height_emu=bio_h,
                )
                tf = tx.text_frame
                tf.word_wrap = True
                first = True
                if bio:
                    p = tf.paragraphs[0]
                    self._set_paragraph(
                        p, text=bio[:280],
                        size_pt=11, bold=False,
                        font_name=fonts["body"],
                        rgb=_hex_to_rgb(palette["text_secondary"]),
                    )
                    first = False
                if linkedin:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    self._set_paragraph(
                        p, text=linkedin,
                        size_pt=10, bold=False,
                        font_name=fonts["mono"] if fonts.get("mono") else fonts["body"],
                        rgb=_hex_to_rgb(palette["text_muted"]),
                    )
                    # Real clickable hyperlink to the LinkedIn profile.
                    if p.runs:
                        self._set_run_hyperlink(p.runs[0], linkedin)

    def _render_image_full_slide(self, pptx_slide, slide, palette, fonts) -> None:
        # Try to embed the real image full-bleed across the slide. On
        # any failure (no URL, fetch error, unsupported mime) we render
        # an honest placeholder showing the URL/prompt verbatim — never
        # a fabricated stock photo.
        url_raw = slide.get("image_url")
        url = str(url_raw).strip() if isinstance(url_raw, str) and url_raw.strip() else ""
        prompt_raw = slide.get("image_prompt")
        prompt = str(prompt_raw).strip() if isinstance(prompt_raw, str) and prompt_raw.strip() else ""
        headline = str(slide.get("headline") or "").strip()
        subheadline = str(slide.get("subheadline") or "").strip()
        # Big image band across the full slide.
        band_top = Inches(0.6).emu
        band_h = _H_EMU - Inches(2.4).emu
        embedded = False
        if url:
            embedded = self._add_picture_safe(
                pptx_slide, url,
                left_emu=_MARGIN_EMU, top_emu=band_top,
                width_emu=_W_EMU - 2 * _MARGIN_EMU, height_emu=band_h,
                link_url=self._first_link_url(slide, {"image"}),
            )
        if not embedded and (url or prompt):
            self._add_text_block(
                pptx_slide,
                text=(f"[Image: {url}]" if url else f"[Image prompt: {prompt[:200]}]"),
                left_emu=_MARGIN_EMU, top_emu=band_top,
                width_emu=_W_EMU - 2 * _MARGIN_EMU, height_emu=band_h,
                size_pt=14, bold=False,
                font_name=fonts["mono"] if fonts.get("mono") else fonts["body"],
                color_hex=palette["text_muted"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
        # Caption (headline + subheadline) at bottom — only when present.
        if headline:
            self._add_text_block(
                pptx_slide, text=headline,
                left_emu=_MARGIN_EMU,
                top_emu=_H_EMU - Inches(1.6).emu,
                width_emu=_W_EMU - 2 * _MARGIN_EMU, height_emu=Inches(0.7).emu,
                size_pt=24, bold=True,
                font_name=fonts["heading"],
                color_hex=palette["text_primary"],
                align=PP_ALIGN.LEFT,
            )
        if subheadline:
            self._add_text_block(
                pptx_slide, text=subheadline,
                left_emu=_MARGIN_EMU,
                top_emu=_H_EMU - Inches(0.9).emu,
                width_emu=_W_EMU - 2 * _MARGIN_EMU, height_emu=Inches(0.5).emu,
                size_pt=14, bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_secondary"],
                align=PP_ALIGN.LEFT,
            )
        # If we have neither image nor headline, fall back to bullets.
        if not (url or prompt or headline or subheadline):
            self._render_bullets_slide(pptx_slide, slide, palette, fonts)

    def _render_failure_slide(
        self, prs, slide, palette, fonts, error_msg: str,
    ) -> None:
        try:
            pptx_slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._set_background(pptx_slide, palette)
            self._add_text_block(
                pptx_slide,
                text=f"Slide {slide.get('index', '?')} could not be exported.",
                left_emu=_MARGIN_EMU,
                top_emu=Inches(2.5).emu,
                width_emu=_W_EMU - 2 * _MARGIN_EMU,
                height_emu=Inches(1.0).emu,
                size_pt=24, bold=True,
                font_name=fonts["heading"],
                color_hex=palette["text_primary"],
                align=PP_ALIGN.CENTER,
            )
            self._add_text_block(
                pptx_slide,
                text=f"Reason: {error_msg}",
                left_emu=_MARGIN_EMU,
                top_emu=Inches(3.7).emu,
                width_emu=_W_EMU - 2 * _MARGIN_EMU,
                height_emu=Inches(1.0).emu,
                size_pt=14, bold=False,
                font_name=fonts["body"],
                color_hex=palette["text_muted"],
                align=PP_ALIGN.CENTER,
            )
        except Exception:  # noqa: BLE001
            pass  # last-resort safety; never raise out of build()

    def _render_link_cta(self, pptx_slide, slide, palette, fonts) -> None:
        links = self._slide_links(slide)
        if not links:
            return
        intent = str(slide.get("intent") or "").lower()
        preferred = next((l for l in links if l.get("target") == "button"), None)
        if preferred is None and intent not in {"ask", "closing", "thanks", "thank_you"}:
            return
        link = preferred or links[0]
        url = link.get("url", "")
        if not url:
            return
        label = (link.get("label") or "Open link")[:42]
        left = _MARGIN_EMU
        top = _H_EMU - Inches(0.85).emu
        width = Inches(2.5).emu
        height = Inches(0.38).emu
        shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(left), Emu(top), Emu(width), Emu(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(palette["accent"])
        shape.line.fill.background()
        self._set_shape_hyperlink(shape, url)
        if shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            tf.margin_left = Inches(0.08).emu
            tf.margin_right = Inches(0.08).emu
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.name = fonts["body"]
            run.font.color.rgb = _hex_to_rgb("#FFFFFF")
            self._set_run_hyperlink(run, url)

    def _render_citation_footer(self, pptx_slide, slide, palette, fonts) -> None:
        # Collect (title, url) pairs from citations[] and enrichment.sources[].
        # We never invent a citation URL; the run.hyperlink is set only when a
        # real URL is present in the DSL.
        pairs: list[tuple[str, str]] = []
        for c in slide.get("citations") or []:
            if not isinstance(c, dict):
                continue
            url = str(c.get("url", "")).strip()
            title = str(c.get("title", "")).strip()
            if url:
                pairs.append((title, url))
            if len(pairs) >= 3:
                break
        for link in self._slide_links(slide)[: max(0, 3 - len(pairs))]:
            url = str(link.get("url", "")).strip()
            title = str(link.get("label", "")).strip()
            if url and (title, url) not in pairs:
                pairs.append((title, url))
        enrichment = slide.get("enrichment") or {}
        for src in (enrichment.get("sources") or [])[: max(0, 3 - len(pairs))]:
            if not isinstance(src, dict):
                continue
            url = str(src.get("url", "")).strip()
            title = str(src.get("title", "")).strip()
            if url:
                pairs.append((title, url))
        if not pairs:
            return

        # Render as a single textbox with multiple runs. Each URL becomes
        # its own run with a real clickable hyperlink. Separators between
        # citations are plain runs with no hyperlink.
        tx = self._add_textbox(
            pptx_slide,
            left_emu=_MARGIN_EMU,
            top_emu=_H_EMU - Inches(0.4).emu,
            width_emu=_W_EMU - 2 * _MARGIN_EMU,
            height_emu=Inches(0.3).emu,
        )
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        def _add_run(text: str, *, url: Optional[str] = None) -> None:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(8)
            run.font.name = fonts["body"]
            run.font.color.rgb = _hex_to_rgb(
                palette["accent"] if url else palette["text_muted"]
            )
            if url:
                self._set_run_hyperlink(run, url)

        _add_run("Sources: ")
        for i, (title, url) in enumerate(pairs):
            if i > 0:
                _add_run("  ·  ")
            label = (title or url)[:160]
            _add_run(label, url=url)


__all__ = ["V4PptxBuilder"]
