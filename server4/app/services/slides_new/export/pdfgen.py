"""
PPTX Export Builder - Phase 2
Generates PowerPoint presentations from slide data.
Uses python-pptx for native PPTX generation.
"""

import io
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from app.services.slides_new.design.system import DesignSystem


class PptxExporter:
    """
    PPTX Export Builder - creates PowerPoint presentations from slide data.

    Features:
    - 16:9 widescreen format
    - Custom layouts per slide type
    - Charts and tables support
    - Image embedding
    - Professional typography
    """

    LAYOUT_MAP = {
        "title-hero": 0,
        "title": 0,
        "two-column": 1,
        "bullets": 1,
        "bullets-with-image": 1,
        "chart": 5,
        "team-grid": 5,
        "comparison": 5,
        "kpi-dashboard": 5,
        "timeline": 5,
        "quote": 6,
        "blank": 6,
    }

    def __init__(self, design_system: Optional[DesignSystem] = None):
        self.design = design_system or DesignSystem()
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    def _create_title_slide(self, title: str, subtitle: str = ""):
        """Create title slide (title-hero layout)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])

        title_shape = slide.shapes.title
        title_shape.text = title

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle

    def _create_content_slide(
        self, title: str, bullets: List[str] = None, layout_idx: int = 1
    ):
        """Create content slide with bullets"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        title_shape = slide.shapes.title
        title_shape.text = title

        if bullets:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

    def _create_two_column_slide(
        self, title: str, left_content: str, right_content: str = ""
    ):
        """Create two column slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = title

        left = Inches(0.5)
        top = Inches(2)
        width = Inches(5.5)
        height = Inches(5)

        txLeft = slide.shapes.add_textbox(left, top, width, height)
        tf = txLeft.text_frame
        tf.text = left_content

    def _create_chart_slide(self, title: str, chart_data: Dict[str, Any]):
        """Create slide with chart"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        slide.shapes.title.text = title

    def _create_quote_slide(self, quote: str, author: str = ""):
        """Create quote slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(2)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f'"{quote}"'
        p.font.size = Pt(28)
        p.font.italic = True

        if author:
            p2 = tf.add_paragraph()
            p2.text = f"- {author}"
            p2.font.size = Pt(18)

    def export_presentation(
        self, slides_data: List[Dict[str, Any]], metadata: Dict[str, Any] = None
    ) -> bytes:
        """
        Export presentation from slides data.

        Args:
            slides_data: List of slide dictionaries with content
            metadata: Presentation metadata

        Returns:
            PPTX file bytes
        """
        for slide_data in slides_data:
            layout = slide_data.get("layout", "bullets")
            title = slide_data.get("title", "Untitled")
            content = slide_data.get("content", {})
            purpose = slide_data.get("purpose", "")

            if layout == "title-hero" or (layout == "title" and len(slides_data) == 1):
                self._create_title_slide(
                    title=title,
                    subtitle=metadata.get("description", "") if metadata else "",
                )
            elif layout == "quote":
                quote_data = content.get("quote", {})
                self._create_quote_slide(
                    quote=quote_data.get("text", ""),
                    author=quote_data.get("author", ""),
                )
            elif layout == "two-column":
                bullets = content.get("bullets", [])
                mid = len(bullets) // 2
                left = "\n".join(bullets[:mid]) if mid > 0 else ""
                right = "\n".join(bullets[mid:]) if mid < len(bullets) else ""
                self._create_two_column_slide(title, left, right)
            else:
                bullets = content.get("bullets", [])
                self._create_content_slide(title, bullets)

        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def export_to_file(self, slides_data: List[Dict], filepath: str):
        """Export to file"""
        pptx_bytes = self.export_presentation(slides_data)
        with open(filepath, "wb") as f:
            f.write(pptx_bytes)


def export_to_pptx(
    presentation: Dict[str, Any], filepath: Optional[str] = None
) -> bytes:
    """
    Convenience function to export presentation to PPTX.

    Args:
        presentation: Full presentation dict with slides and metadata
        filepath: Optional filepath to save to

    Returns:
        PPTX bytes
    """
    preset = presentation.get("design_system", {}).get("preset", "yc_pitch")
    design = DesignSystem(preset)

    exporter = PptxExporter(design)
    slides = presentation.get("slides", [])
    metadata = presentation.get("metadata", {})

    pptx_bytes = exporter.export_presentation(slides, metadata)

    if filepath:
        with open(filepath, "wb") as f:
            f.write(pptx_bytes)

    return pptx_bytes
