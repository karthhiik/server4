"""Phase 13 — V4 PPTX export unit tests.

Verify each layout produces correct PPTX shapes by reading the produced
file back with python-pptx. No fake-data invariants are also asserted
(empty fields → omitted shapes; never placeholder copy).
"""

from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from app.services.v4.pptx_export import (
    V4PptxBuilder,
    _coerce_float,
    _hex_to_rgb,
)


# ── fixtures ──────────────────────────────────────────────────────


_TOKENS = {
    "palette": {
        "primary": "#0043CE",
        "secondary": "#1E40AF",
        "accent": "#7C3AED",
        "background": "#FFFFFF",
        "surface": "#F8FAFC",
        "text_primary": "#161616",
        "text_secondary": "#525252",
        "text_muted": "#6F6F6F",
        "success": "#16A34A",
        "warning": "#F59E0B",
        "danger": "#DC2626",
        "chart": ["#0043CE", "#7C3AED", "#16A34A", "#F59E0B"],
    },
    "fonts": {"heading": "Inter", "body": "Inter"},
}


def _open(b: bytes) -> Presentation:
    return Presentation(io.BytesIO(b))


def _all_text(slide) -> str:
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                out.append(para.text)
    return "\n".join(out)


# ── construction & guards ─────────────────────────────────────────


def test_build_empty_raises():
    with pytest.raises(ValueError):
        V4PptxBuilder().build([])


def test_build_filters_non_dict_entries_then_errors_when_empty():
    with pytest.raises(ValueError):
        V4PptxBuilder().build(["not a dict", 42, None])  # type: ignore[list-item]


def test_hex_to_rgb_handles_short_form_and_garbage():
    rgb = _hex_to_rgb("#fff")
    assert (rgb[0], rgb[1], rgb[2]) == (255, 255, 255)
    rgb2 = _hex_to_rgb("not-a-color")
    assert (rgb2[0], rgb2[1], rgb2[2]) == (0, 0, 0)


def test_coerce_float_strips_currency_percent_suffix():
    assert _coerce_float("$4.1M") == pytest.approx(4_100_000.0)
    assert _coerce_float("18%") == 18.0
    assert _coerce_float("1,200") == 1200.0
    assert _coerce_float(None) is None
    assert _coerce_float("garbage") is None


# ── layout coverage ───────────────────────────────────────────────


def test_title_slide_has_headline_and_subhead():
    slide = {
        "index": 0, "intent": "title", "layout": "title",
        "headline": "Acme · Investor Deck",
        "subheadline": "Series A — Q1 2026",
        "body": "Confidential",
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    assert len(prs.slides) == 1
    text = _all_text(prs.slides[0])
    assert "Acme · Investor Deck" in text
    assert "Series A" in text
    assert "Confidential" in text


def test_bullets_slide_emits_one_paragraph_per_bullet():
    slide = {
        "index": 0, "intent": "problem", "layout": "bullets",
        "headline": "Customer pain",
        "bullets": ["Slow imports", "Manual reformatting", "No version control"],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    for b in slide["bullets"]:
        assert b in text
    # Bullets are prefixed with •
    assert text.count("•") >= 3


def test_stat_grid_slide_renders_value_and_label():
    slide = {
        "index": 0, "intent": "traction", "layout": "stat-grid",
        "headline": "Traction",
        "stat_blocks": [
            {"value": "12K", "label": "Monthly active users"},
            {"value": "$420K", "label": "ARR"},
            {"value": "3.4x", "label": "YoY growth"},
        ],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    for s in slide["stat_blocks"]:
        assert s["value"] in text
        assert s["label"] in text


def test_chart_slide_creates_native_pptx_chart():
    slide = {
        "index": 0, "intent": "market", "layout": "chart",
        "headline": "TAM",
        "chart": {
            "type": "bar",
            "series_label": "USD (M)",
            "data": [
                {"label": "2023", "value": 10},
                {"label": "2024", "value": 18},
                {"label": "2025", "value": 32},
            ],
        },
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    pptx_slide = prs.slides[0]
    chart_shapes = [s for s in pptx_slide.shapes if s.has_chart]
    assert len(chart_shapes) == 1
    assert chart_shapes[0].chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED


def test_chart_slide_drops_invalid_points_no_fake_zero():
    """A row with a non-coercible value must be DROPPED, not forced to 0."""
    slide = {
        "index": 0, "intent": "market", "layout": "chart",
        "headline": "TAM",
        "chart": {
            "type": "bar",
            "data": [
                {"label": "2023", "value": "garbage"},
                {"label": "2024", "value": 18},
            ],
        },
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    pptx_slide = prs.slides[0]
    chart_shapes = [s for s in pptx_slide.shapes if s.has_chart]
    assert len(chart_shapes) == 1
    # Only one valid category survived.
    cats = list(chart_shapes[0].chart.plots[0].categories)
    assert cats == ["2024"]


def test_chart_slide_with_no_usable_data_falls_back_to_bullets_no_fake_chart():
    slide = {
        "index": 0, "intent": "market", "layout": "chart",
        "headline": "Headline only",
        "bullets": ["fallback bullet"],
        "chart": {"type": "bar", "data": [{"label": "x", "value": "garbage"}]},
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    pptx_slide = prs.slides[0]
    chart_shapes = [s for s in pptx_slide.shapes if s.has_chart]
    assert len(chart_shapes) == 0  # no fabricated chart
    assert "fallback bullet" in _all_text(pptx_slide)


def test_table_slide_emits_table_shape_with_correct_dimensions():
    slide = {
        "index": 0, "intent": "competition", "layout": "table",
        "headline": "Competitive matrix",
        "table": {
            "headers": ["Feature", "Acme", "CompetitorX"],
            "rows": [
                ["Real-time edits", "Yes", "No"],
                ["AI generation", "Native", "Plugin"],
            ],
        },
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    table_shapes = [s for s in prs.slides[0].shapes if s.has_table]
    assert len(table_shapes) == 1
    tbl = table_shapes[0].table
    assert len(tbl.columns) == 3
    assert len(tbl.rows) == 3   # 1 header + 2 data
    assert tbl.cell(0, 0).text == "Feature"
    assert tbl.cell(2, 2).text == "Plugin"


def test_quote_slide_includes_attribution():
    slide = {
        "index": 0, "intent": "vision", "layout": "quote",
        "quote": {"text": "Move fast and ship.", "attribution": "Founder"},
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    assert "Move fast and ship." in text
    assert "Founder" in text


def test_speaker_notes_embedded_in_notes_slide():
    slide = {
        "index": 0, "intent": "title", "layout": "title",
        "headline": "Hi",
        "speaker_notes": "Open with a strong hook about the market shift.",
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    notes_text = prs.slides[0].notes_slide.notes_text_frame.text
    assert "strong hook" in notes_text


def test_empty_speaker_notes_does_not_emit_placeholder():
    """No-fake-data: empty notes must NOT put a placeholder string."""
    slide = {
        "index": 0, "intent": "title", "layout": "title",
        "headline": "Hi", "speaker_notes": "   ",
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    notes_text = prs.slides[0].notes_slide.notes_text_frame.text
    # python-pptx initializes notes to empty string. We must not write
    # "Add notes" or any other placeholder copy.
    assert "Add notes" not in notes_text
    assert "TODO" not in notes_text
    assert "placeholder" not in notes_text.lower()


def test_unresolved_team_slide_exports_no_fake_members():
    slide = {
        "index": 0,
        "intent": "team",
        "layout": "team-grid",
        "headline": "Team",
        "subheadline": "Verified team details are pending",
        "requires_user_input": True,
        "team_members": [],
        "bullets": ["Founder & CEO", "Co-Founder & CTO"],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])

    assert "Team" in text
    assert "Founder & CEO" not in text
    assert "Co-Founder & CTO" not in text


def test_unknown_layout_falls_back_to_bullets_default():
    slide = {
        "index": 0, "intent": "custom", "layout": "totally-made-up-layout",
        "headline": "Headline",
        "bullets": ["one", "two"],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    assert "Headline" in text
    assert "one" in text and "two" in text


def test_image_url_renders_honest_placeholder_not_blank():
    """No-fake-data: missing image renders an [Image: <url>] caption,
    NOT a fabricated diagram or generic stock filename."""
    slide = {
        "index": 0, "intent": "solution", "layout": "image-right",
        "headline": "Workflow",
        "bullets": ["step 1"],
        "image_url": "https://example.com/diagram.png",
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    assert "[Image:" in text
    assert "example.com/diagram.png" in text


def test_citation_footer_only_appears_when_citations_exist():
    no_cite = {
        "index": 0, "intent": "market", "layout": "bullets",
        "headline": "X", "bullets": ["a"],
    }
    cite = {
        "index": 0, "intent": "market", "layout": "bullets",
        "headline": "X", "bullets": ["a"],
        "citations": [{"title": "Gartner 2025", "url": "https://gartner.com/x"}],
    }
    prs_a = _open(V4PptxBuilder().build([no_cite], _TOKENS))
    prs_b = _open(V4PptxBuilder().build([cite], _TOKENS))
    text_a = _all_text(prs_a.slides[0])
    text_b = _all_text(prs_b.slides[0])
    assert "Sources:" not in text_a
    assert "Sources:" in text_b
    assert "Gartner 2025" in text_b
    # 100% fidelity: the citation URL is attached as a real clickable
    # hyperlink to the run rather than dumped as plain text.
    hyperlinks = []
    for shape in prs_b.slides[0].shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                addr = getattr(run.hyperlink, "address", None)
                if addr:
                    hyperlinks.append(addr)
    assert "https://gartner.com/x" in hyperlinks


def test_all_allowed_layouts_render_without_crash():
    """Every layout in v4 _ALLOWED_LAYOUTS must produce a valid slide."""
    layouts = [
        "title", "two-column", "three-column", "bullets", "stat-grid",
        "quote", "chart", "table", "timeline", "comparison", "diagram",
        "image-full", "image-left", "image-right", "team",
    ]
    slides = []
    for i, lay in enumerate(layouts):
        slides.append({
            "index": i,
            "intent": "generic",
            "layout": lay,
            "headline": f"Slide {i}: {lay}",
            "subheadline": "sub",
            "bullets": ["a", "b", "c"],
            "stat_blocks": [{"value": "1", "label": "x"}],
            "quote": {"text": "Q", "attribution": "Z"},
            "chart": {"type": "bar", "data": [{"label": "a", "value": 1}]},
            "table": {"headers": ["h"], "rows": [["v"]]},
            "timeline": {"events": [{"date": "2025", "title": "T", "description": "D"}]},
            "comparison": {"columns": [{"title": "Us", "items": ["x"]}, {"title": "Them", "items": ["y"]}]},
            "diagram": {"nodes": [{"label": "N1"}, {"label": "N2"}]},
        })
    prs = _open(V4PptxBuilder().build(slides, _TOKENS))
    assert len(prs.slides) == len(layouts)


def test_design_tokens_palette_drives_background_color():
    custom = {
        "palette": {**_TOKENS["palette"], "background": "#101820"},
        "fonts": _TOKENS["fonts"],
    }
    slide = {"index": 0, "intent": "title", "layout": "title", "headline": "X"}
    prs = _open(V4PptxBuilder().build([slide], custom))
    bg = prs.slides[0].background.fill
    assert str(bg.fore_color.rgb) == "101820"


def test_metadata_title_set_on_core_properties():
    slide = {"index": 0, "intent": "title", "layout": "title", "headline": "X"}
    prs = _open(V4PptxBuilder().build(
        [slide], _TOKENS, metadata={"title": "My Deck", "company": "Acme"},
    ))
    assert prs.core_properties.title == "My Deck"


def test_slide_failure_renders_error_slide_not_silent_drop():
    """A slide that triggers an internal exception still produces output
    so the user knows something went wrong — but contains NO fabricated
    content."""
    builder = V4PptxBuilder()
    # Patch _render_slide to raise on the second slide.
    real = builder._render_slide
    calls = {"n": 0}

    def boom(prs, slide, palette, fonts):
        calls["n"] += 1
        if calls["n"] == 2:
            raise RuntimeError("synthetic")
        return real(prs, slide, palette, fonts)

    builder._render_slide = boom  # type: ignore[method-assign]
    out = builder.build(
        [
            {"index": 0, "intent": "title", "layout": "title", "headline": "OK"},
            {"index": 1, "intent": "title", "layout": "title", "headline": "BOOM"},
        ],
        _TOKENS,
    )
    prs = _open(out)
    assert len(prs.slides) == 2
    text2 = _all_text(prs.slides[1])
    assert "could not be exported" in text2
    # The original (potentially partial) headline must not be guessed at
    # — we don't claim it succeeded.
    assert "BOOM" not in text2


def test_no_design_tokens_uses_safe_defaults():
    slide = {"index": 0, "intent": "title", "layout": "title", "headline": "X"}
    out = V4PptxBuilder().build([slide])  # no tokens
    prs = _open(out)
    assert len(prs.slides) == 1


def test_slide_ordering_preserved_by_index():
    slides = [
        {"index": 2, "intent": "title", "layout": "title", "headline": "C"},
        {"index": 0, "intent": "title", "layout": "title", "headline": "A"},
        {"index": 1, "intent": "title", "layout": "title", "headline": "B"},
    ]
    prs = _open(V4PptxBuilder().build(slides, _TOKENS))
    texts = [_all_text(s) for s in prs.slides]
    assert "A" in texts[0]
    assert "B" in texts[1]
    assert "C" in texts[2]


# ── premium-mode layouts: team + image-full ───────────────────────


def test_team_slide_renders_each_member_name_and_role():
    slide = {
        "index": 0,
        "intent": "team",
        "layout": "team",
        "headline": "Founding Team",
        "team_members": [
            {"name": "Ada Lovelace", "role": "CEO", "bio": "Mathematician.",
             "linkedin_url": "https://linkedin.com/in/ada"},
            {"name": "Grace Hopper", "role": "CTO", "bio": "Compiler pioneer."},
            {"name": "Alan Turing", "role": "Chief Scientist"},
        ],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    assert "Founding Team" in text
    assert "Ada Lovelace" in text and "CEO" in text
    assert "Grace Hopper" in text and "CTO" in text
    assert "Alan Turing" in text and "Chief Scientist" in text
    # Bio + linkedin only when present — never invented
    assert "Mathematician." in text
    assert "https://linkedin.com/in/ada" in text


def test_team_slide_filters_invalid_entries_no_fabricated_names():
    slide = {
        "index": 0,
        "intent": "team",
        "layout": "team",
        "headline": "Team",
        "team_members": [
            {"name": "Real Person", "role": "Founder"},
            {"name": "", "role": "Empty"},          # filtered: blank name
            {"role": "No Name"},                     # filtered: missing name
            "not a dict",                            # filtered: wrong type
            {"name": "  ", "role": "Whitespace"},   # filtered: whitespace name
        ],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    assert "Real Person" in text
    # None of the dropped roles should leak in
    assert "Empty" not in text
    assert "No Name" not in text
    assert "Whitespace" not in text


def test_team_slide_with_no_valid_members_falls_back_to_bullets():
    slide = {
        "index": 0,
        "intent": "team",
        "layout": "team",
        "headline": "Team",
        "bullets": ["We are hiring", "Roles open"],
        "team_members": [{"role": "x"}, "garbage"],  # all invalid
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    # Falls through cleanly to bullets — never fabricates a member
    assert "We are hiring" in text
    assert "Roles open" in text


def test_image_full_slide_renders_url_placeholder_and_caption():
    slide = {
        "index": 0,
        "intent": "vision",
        "layout": "image-full",
        "headline": "The Future Is Now",
        "subheadline": "AI-native presentations",
        "image_url": "https://cdn.example.com/hero.jpg",
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    # Honest placeholder — image bytes are never fetched at export
    assert "[Image: https://cdn.example.com/hero.jpg]" in text
    assert "The Future Is Now" in text
    assert "AI-native presentations" in text


def test_image_full_slide_uses_prompt_when_url_missing():
    slide = {
        "index": 0,
        "intent": "vision",
        "layout": "image-full",
        "headline": "Vision",
        "image_prompt": "Cinematic shot of a city skyline at dusk",
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    assert "[Image prompt: Cinematic shot of a city skyline at dusk]" in text
    assert "Vision" in text


def test_image_full_slide_with_nothing_falls_back_to_bullets_no_blank_slide():
    slide = {
        "index": 0,
        "intent": "x",
        "layout": "image-full",
        "bullets": ["Only bullets here"],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    text = _all_text(prs.slides[0])
    # Empty image-full slide must not produce a blank rectangle —
    # falls back to bullets renderer.
    assert "Only bullets here" in text
    assert "[Image" not in text  # no fake placeholder for missing image


# ── 100 %-fidelity feature coverage ────────────────────────────────
#
# These tests verify the upgrade from "honest text placeholder" rendering
# to real PPTX shapes: embedded pictures, rounded card backgrounds,
# arrow connectors, clickable hyperlinks, and page-number footers.
# All assertions go through python-pptx's read-back path so we are
# verifying actual XML, not in-memory state.

# 67-byte 1×1 transparent PNG. Real bytes — not fabricated.
_PNG_1PX = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00"
    b"\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\rIDATx"
    b"\x9cc\x00\x01\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND"
    b"\xaeB`\x82"
)


def _count_pictures(slide) -> int:
    n = 0
    for shape in slide.shapes:
        # MSO_SHAPE_TYPE.PICTURE = 13
        if getattr(shape, "shape_type", None) == 13:
            n += 1
    return n


def _count_shapes_of_type(slide, autoshape_type: int) -> int:
    """Count auto-shapes whose `auto_shape_type` enum int matches."""
    n = 0
    for shape in slide.shapes:
        ast = None
        try:
            ast = shape.auto_shape_type
        except Exception:  # noqa: BLE001
            ast = None
        if ast is not None and int(ast) == autoshape_type:
            n += 1
    return n


def test_image_resolver_embeds_real_picture_in_image_full():
    """Resolver injection bypasses HTTP and embeds real image bytes."""
    calls: list[str] = []

    def _resolver(url: str) -> bytes | None:
        calls.append(url)
        return _PNG_1PX

    builder = V4PptxBuilder(image_resolver=_resolver)
    slide = {
        "index": 0,
        "intent": "x",
        "layout": "image-full",
        "image_url": "https://example.com/hero.png",
        "headline": "Hero",
    }
    prs = _open(builder.build([slide], _TOKENS))
    assert calls == ["https://example.com/hero.png"]
    assert _count_pictures(prs.slides[0]) == 1
    text = _all_text(prs.slides[0])
    # No "[Image: ...]" text placeholder should appear once the picture
    # was embedded successfully.
    assert "[Image:" not in text


def test_image_resolver_returning_none_falls_back_to_text_no_fake_image():
    """When resolver yields no bytes, render honest text placeholder."""
    builder = V4PptxBuilder(image_resolver=lambda _u: None)
    slide = {
        "index": 0, "intent": "x", "layout": "image-full",
        "image_url": "https://example.com/missing.png",
    }
    prs = _open(builder.build([slide], _TOKENS))
    assert _count_pictures(prs.slides[0]) == 0
    assert "[Image:" in _all_text(prs.slides[0])


def test_image_resolver_used_for_image_left_image_right_team_title():
    """Real image embedding works on every layout that has an image."""
    builder = V4PptxBuilder(image_resolver=lambda _u: _PNG_1PX)
    slides = [
        {
            "index": 0, "intent": "x", "layout": "title",
            "headline": "T", "company_icon_url": "https://e.com/icon.png",
        },
        {
            "index": 1, "intent": "x", "layout": "image-left",
            "headline": "L", "image_url": "https://e.com/l.png",
            "bullets": ["a"],
        },
        {
            "index": 2, "intent": "x", "layout": "image-right",
            "headline": "R", "image_url": "https://e.com/r.png",
            "bullets": ["a"],
        },
        {
            "index": 3, "intent": "x", "layout": "team",
            "headline": "Team",
            "team_members": [
                {"name": "Alice", "role": "CEO",
                 "photo_url": "https://e.com/a.png"},
            ],
        },
    ]
    prs = _open(builder.build(slides, _TOKENS))
    assert _count_pictures(prs.slides[0]) == 1  # title icon
    assert _count_pictures(prs.slides[1]) == 1  # image-left
    assert _count_pictures(prs.slides[2]) == 1  # image-right
    assert _count_pictures(prs.slides[3]) == 1  # team photo


def test_decode_data_uri_accepts_supported_mime_rejects_others():
    """`_decode_data_uri` must enforce the PPTX-supported MIME whitelist."""
    from app.services.v4.pptx_export import _decode_data_uri

    import base64
    b64 = base64.b64encode(_PNG_1PX).decode("ascii")
    ok = _decode_data_uri(f"data:image/png;base64,{b64}")
    assert ok is not None
    body, mime = ok
    assert body == _PNG_1PX
    assert mime == "image/png"

    # SVG is not embeddable in PPTX → must be rejected.
    svg_b64 = base64.b64encode(b"<svg/>").decode("ascii")
    assert _decode_data_uri(f"data:image/svg+xml;base64,{svg_b64}") is None
    # Malformed → None, never raises.
    assert _decode_data_uri("data:image/png;base64,!!!not-base64!!!") is None
    assert _decode_data_uri("not a data uri at all") is None


def test_data_uri_image_embeds_directly_without_resolver():
    """Photo URLs as data: URIs embed without needing a resolver."""
    import base64
    b64 = base64.b64encode(_PNG_1PX).decode("ascii")
    data_uri = f"data:image/png;base64,{b64}"
    slide = {
        "index": 0, "intent": "x", "layout": "team",
        "headline": "Team",
        "team_members": [
            {"name": "Alice", "role": "CEO", "photo_url": data_uri},
        ],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    assert _count_pictures(prs.slides[0]) == 1


def test_embed_images_false_skips_all_image_fetches():
    """`embed_images=False` produces zero pictures and no resolver calls."""
    calls: list[str] = []

    def _resolver(url: str) -> bytes | None:
        calls.append(url)
        return _PNG_1PX

    builder = V4PptxBuilder(image_resolver=_resolver)
    slide = {
        "index": 0, "intent": "x", "layout": "image-full",
        "image_url": "https://example.com/hero.png",
    }
    prs = _open(builder.build([slide], _TOKENS, embed_images=False))
    assert calls == []
    assert _count_pictures(prs.slides[0]) == 0


def test_stat_grid_renders_rounded_card_backgrounds():
    """Each stat in stat-grid gets a real ROUNDED_RECTANGLE card."""
    from pptx.enum.shapes import MSO_SHAPE
    slide = {
        "index": 0, "intent": "x", "layout": "stat-grid",
        "headline": "Stats",
        "stat_blocks": [
            {"value": "1", "label": "a"},
            {"value": "2", "label": "b"},
            {"value": "3", "label": "c"},
        ],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    cards = _count_shapes_of_type(prs.slides[0], int(MSO_SHAPE.ROUNDED_RECTANGLE))
    assert cards == 3


def test_comparison_highlight_column_gets_card_with_accent_border():
    """Highlighted comparison columns render with a real card behind them."""
    from pptx.enum.shapes import MSO_SHAPE
    slide = {
        "index": 0, "intent": "x", "layout": "comparison",
        "headline": "Vs",
        "comparison": {
            "columns": [
                {"title": "Us", "items": ["a"], "highlight": True},
                {"title": "Them", "items": ["b"]},
            ],
        },
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    # Both columns get a card; the highlighted one has a visible border.
    cards = _count_shapes_of_type(prs.slides[0], int(MSO_SHAPE.ROUNDED_RECTANGLE))
    assert cards == 2


def test_diagram_adjacent_edge_renders_real_right_arrow():
    """Adjacent edges get a real RIGHT_ARROW shape, not just text."""
    from pptx.enum.shapes import MSO_SHAPE
    slide = {
        "index": 0, "intent": "x", "layout": "diagram",
        "headline": "Flow",
        "diagram": {
            "nodes": [{"label": "A"}, {"label": "B"}, {"label": "C"}],
            "edges": [{"from": 0, "to": 1}, {"from": 1, "to": 2}],
        },
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    arrows = _count_shapes_of_type(prs.slides[0], int(MSO_SHAPE.RIGHT_ARROW))
    assert arrows == 2
    # Both edges adjacent → caption "Flow:" should NOT appear.
    assert "Flow:" not in _all_text(prs.slides[0])


def test_diagram_non_adjacent_edge_falls_back_to_caption_no_fake_arrow():
    """Non-adjacent edges must NOT fabricate a rendered arrow."""
    from pptx.enum.shapes import MSO_SHAPE
    slide = {
        "index": 0, "intent": "x", "layout": "diagram",
        "headline": "Flow",
        "diagram": {
            "nodes": [{"label": "A"}, {"label": "B"}, {"label": "C"}],
            "edges": [{"from": 0, "to": 2, "label": "skip"}],
        },
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    arrows = _count_shapes_of_type(prs.slides[0], int(MSO_SHAPE.RIGHT_ARROW))
    assert arrows == 0
    text = _all_text(prs.slides[0])
    assert "Flow:" in text
    assert "A" in text and "C" in text


def test_diagram_resolves_edges_by_node_id_string():
    """Edges using string `id` references resolve correctly."""
    from pptx.enum.shapes import MSO_SHAPE
    slide = {
        "index": 0, "intent": "x", "layout": "diagram",
        "headline": "Flow",
        "diagram": {
            "nodes": [
                {"id": "start", "label": "Start"},
                {"id": "end", "label": "End"},
            ],
            "edges": [{"from": "start", "to": "end"}],
        },
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    arrows = _count_shapes_of_type(prs.slides[0], int(MSO_SHAPE.RIGHT_ARROW))
    assert arrows == 1


def test_page_numbers_added_when_multi_slide_default():
    """`i / N` footer appears on every slide in a multi-slide deck."""
    slides = [
        {"index": 0, "intent": "x", "layout": "bullets",
         "headline": "A", "bullets": ["x"]},
        {"index": 1, "intent": "x", "layout": "bullets",
         "headline": "B", "bullets": ["y"]},
        {"index": 2, "intent": "x", "layout": "bullets",
         "headline": "C", "bullets": ["z"]},
    ]
    prs = _open(V4PptxBuilder().build(slides, _TOKENS))
    for i, sl in enumerate(prs.slides):
        assert f"{i + 1} / 3" in _all_text(sl)


def test_page_numbers_omitted_for_single_slide_no_fake_chrome():
    """Single-slide decks must not render a `1 / 1` footer."""
    slide = {
        "index": 0, "intent": "x", "layout": "bullets",
        "headline": "Solo", "bullets": ["only"],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    assert "1 / 1" not in _all_text(prs.slides[0])


def test_page_numbers_can_be_disabled_via_kwarg():
    """`show_page_numbers=False` suppresses footers even on multi-slide."""
    slides = [
        {"index": 0, "intent": "x", "layout": "bullets",
         "headline": "A", "bullets": ["x"]},
        {"index": 1, "intent": "x", "layout": "bullets",
         "headline": "B", "bullets": ["y"]},
    ]
    prs = _open(
        V4PptxBuilder().build(slides, _TOKENS, show_page_numbers=False)
    )
    for sl in prs.slides:
        assert "1 / 2" not in _all_text(sl)
        assert "2 / 2" not in _all_text(sl)


def test_chart_focus_canonical_layout_routes_to_chart_renderer():
    """Bug 2 regression: the canonical planner layout `chart-focus`
    must route to the native PPTX chart renderer, NOT silently fall
    through to the bullets renderer.
    """
    slide = {
        "index": 0, "intent": "market", "layout": "chart-focus",
        "headline": "TAM growth",
        "chart": {
            "type": "bar",
            "data": [
                {"label": "2024", "value": 12},
                {"label": "2025", "value": 22},
            ],
        },
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    chart_shapes = [s for s in prs.slides[0].shapes if s.has_chart]
    assert len(chart_shapes) == 1, (
        "chart-focus layout fell through to a non-chart renderer"
    )


def test_process_canonical_layout_routes_to_diagram_renderer():
    """Bug 2 regression: the canonical planner layout `process` must
    route through the diagram path so step-flow content keeps its
    directional treatment instead of becoming a generic bullet wall.
    """
    slide = {
        "index": 0, "intent": "how_it_works", "layout": "process",
        "headline": "Onboarding flow",
        "diagram": {
            "nodes": [
                {"id": "a", "label": "Sign up"},
                {"id": "b", "label": "Connect data"},
                {"id": "c", "label": "Insights"},
            ],
            "edges": [{"from": "a", "to": "b"}, {"from": "b", "to": "c"}],
        },
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    from pptx.enum.shapes import MSO_SHAPE
    arrows = _count_shapes_of_type(prs.slides[0], int(MSO_SHAPE.RIGHT_ARROW))
    assert arrows == 2, (
        f"process layout did not render directional arrows (got {arrows})"
    )


def test_team_member_linkedin_renders_clickable_hyperlink():
    """Team-member LinkedIn URLs are real clickable hyperlinks."""
    slide = {
        "index": 0, "intent": "x", "layout": "team",
        "headline": "Team",
        "team_members": [
            {"name": "Alice", "role": "CEO", "bio": "Built X.",
             "linkedin_url": "https://linkedin.com/in/alice"},
        ],
    }
    prs = _open(V4PptxBuilder().build([slide], _TOKENS))
    found = False
    for shape in prs.slides[0].shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if getattr(run.hyperlink, "address", None) == \
                        "https://linkedin.com/in/alice":
                    found = True
    assert found
